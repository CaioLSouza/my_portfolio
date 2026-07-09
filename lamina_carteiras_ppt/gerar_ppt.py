#!/usr/bin/env python3
"""Gera a lâmina PPT da carteira a partir da planilha Excel de gráficos.

O script atualiza, preservando toda a formatação do template:

* **Gráficos** — os gráficos do template guardam no XML as referências às
  células de origem (ex. 'Charts Base 100'!$M$4:$M$10000); o script lê os
  valores atuais dessas células e substitui os dados de cada série.
* **Tabelas de desempenho** (slide "Desempenho") — retornos mensais,
  desempenho por ativo e indicadores (Sharpe, Volatilidade, Beta), lidos
  das abas correspondentes da planilha de gráficos. A tabela de ativos
  cresce ou encolhe conforme o número de linhas na planilha.
* **Tabela de composição** (slide 1) — reconstruída a partir da planilha de
  composição (composicao_*.xlsx), inclusive os agrupamentos mesclados de
  segmento/setor, rating e preço-alvo.
* **Tabela de teses** (slide 4) — peso, recomendação, preço-alvo e o link
  para a tese (montado por ticker). Os comentários, editoriais, são
  preservados casando pelo ticker; papéis novos ficam em branco.
* **Datas** — os textos "D de mês de AAAA" e o título "Mês AAAA" são
  trocados pela data passada em --data (padrão: hoje).

Só continuam manuais: a manchete, o comentário da carteira e o texto dos
comentários por tese (slide 4).

Uso
---
    python gerar_ppt.py --template "Carteira Top Ações - Julho 2026.pptx" \
                        --planilha "Charts Lâmina Carteiras.xlsm" \
                        --saida "Carteira Top Ações - Agosto 2026.pptx" \
                        --carteira top_acoes --data 04/08/2026

Requisitos: pip install python-pptx openpyxl
"""

import argparse
import datetime as dt
import os
import re
import sys
from copy import deepcopy

import openpyxl
from openpyxl.utils import range_boundaries
from pptx import Presentation
from pptx.chart.data import CategoryChartData

# Namespaces do DrawingML (chartSpace e "main", usado em texto/tabelas)
C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Link para a tese de cada papel (o final é o ticker)
LINK_TESE = "https://conteudos.xpi.com.br/acoes/{ticker}"

MESES_ABREV = ["jan", "fev", "mar", "abr", "mai", "jun",
               "jul", "ago", "set", "out", "nov", "dez"]
MESES_EXTENSO = ["janeiro", "fevereiro", "março", "abril", "maio", "junho",
                 "julho", "agosto", "setembro", "outubro", "novembro", "dezembro"]

# Abas de origem das tabelas, por carteira. "composicao" é o nome do arquivo
# gerado na pasta output/ (planilha à parte, com abas PT/PT_data).
CARTEIRAS = {
    "top_acoes": {
        "metricas": "Top Ações Portfolio Metrics",
        "retornos": "performance_top_ações",
        "ativos": "Desempenho ativo top ações_last",
        "composicao": "composicao_top_acoes.xlsx",
    },
    "top_div": {
        "metricas": "Top Div Portfolio Metrics",
        "retornos": "performance_top_divs",
        "ativos": "Desempenho ativo div_last",
        "composicao": "composicao_top_div.xlsx",
    },
    "top_smll": {
        "metricas": "Top SMLL Portfolio Metrics",
        "retornos": "performance_top_smll",
        "ativos": "Desempenho ativo small_last",
        "composicao": "composicao_top_smll.xlsx",
    },
}


def qn(tag: str) -> str:
    return f"{{{C_NS}}}{tag}"


def qnA(tag: str) -> str:
    return f"{{{A_NS}}}{tag}"


def mes_ano(d) -> str:
    """2025-12-01 -> 'dez-25'"""
    return f"{MESES_ABREV[d.month - 1]}-{d.year % 100:02d}"


def pct(v, casas=1) -> str:
    return "" if v is None else f"{v:.{casas}%}"


def pct_br(v, casas=1) -> str:
    """Fração -> percentual no padrão brasileiro: 0.05 -> '5,0%'."""
    return "" if v is None else f"{v * 100:.{casas}f}%".replace(".", ",")


def rs_br(v) -> str:
    """Número -> moeda no padrão brasileiro: 63 -> 'R$ 63,00'."""
    if v is None:
        return ""
    s = f"{v:,.2f}"  # 1,234.56 (separadores americanos)
    return "R$ " + s.replace(",", "\x00").replace(".", ",").replace("\x00", ".")


# ---------------------------------------------------------------------------
# Gráficos
# ---------------------------------------------------------------------------

def parse_ref(ref: str):
    """Divide uma referência tipo "'Charts Base 100'!$M$4:$M$10000" em
    (nome_da_aba, range)."""
    m = re.match(r"^'((?:[^']|'')+)'!(.+)$", ref)
    if m:
        return m.group(1).replace("''", "'"), m.group(2)
    sheet, _, rng = ref.partition("!")
    return sheet, rng


def read_range(wb, ref: str):
    """Lê uma referência de coluna única da planilha e devolve a lista de
    valores (sem cortar vazios — o corte é feito depois, pela categoria)."""
    sheet_name, rng = parse_ref(ref)
    if sheet_name not in wb.sheetnames:
        raise KeyError(f"Aba '{sheet_name}' não encontrada na planilha")
    ws = wb[sheet_name]
    min_col, min_row, max_col, max_row = range_boundaries(rng.replace("$", ""))
    max_row = min(max_row, ws.max_row)
    if min_col != max_col:
        raise ValueError(f"Referência {ref} não é de coluna única")
    return [ws.cell(row=r, column=min_col).value for r in range(min_row, max_row + 1)]


def read_cell(wb, ref: str):
    sheet_name, rng = parse_ref(ref)
    if sheet_name not in wb.sheetnames:
        return None
    return wb[sheet_name][rng.replace("$", "").split(":")[0]].value


def series_info(ser_el):
    """Extrai de um <c:ser> as referências de categoria/valores e o nome."""
    def f_ref(parent_tag, ref_tag):
        parent = ser_el.find(qn(parent_tag))
        if parent is None:
            return None
        el = parent.find(f"{qn(ref_tag)}/{qn('f')}")
        return el.text if el is not None else None

    cat_ref = f_ref("cat", "numRef") or f_ref("cat", "strRef")
    val_ref = f_ref("val", "numRef")
    name_ref = f_ref("tx", "strRef")
    # nome do gráfico: texto literal (<c:tx><c:v>) ou cache da referência
    cached = ser_el.find(f"{qn('tx')}/{qn('v')}")
    if cached is None:
        cached = ser_el.find(f"{qn('tx')}/{qn('strRef')}/{qn('strCache')}/{qn('pt')}/{qn('v')}")
    cached_name = cached.text if cached is not None else None
    return cat_ref, val_ref, name_ref, cached_name


def update_chart(chart, wb, label: str) -> bool:
    """Reconstrói os dados de um gráfico a partir da planilha. Devolve True
    se o gráfico foi atualizado."""
    sers = chart._chartSpace.findall(f".//{qn('ser')}")
    if not sers:
        print(f"  [{label}] sem séries — ignorado")
        return False

    cat_ref, _, _, _ = series_info(sers[0])
    if cat_ref is None:
        print(f"  [{label}] sem referência de categorias — ignorado")
        return False

    raw_cats = read_range(wb, cat_ref)
    # Os ranges do template vão até a linha 10000; corta no fim real dos dados
    n = len(raw_cats)
    while n > 0 and raw_cats[n - 1] is None:
        n -= 1
    cats = raw_cats[:n]
    if not cats:
        print(f"  [{label}] categorias vazias em {cat_ref} — ignorado")
        return False
    if isinstance(cats[0], dt.datetime):
        cats = [c.date() if isinstance(c, dt.datetime) else c for c in cats]

    cat_fmt = "dd/mm/yyyy" if isinstance(cats[0], dt.date) else "General"
    data = CategoryChartData(number_format=cat_fmt)
    data.categories = cats

    for ser_el in sers:
        _, val_ref, name_ref, cached_name = series_info(ser_el)
        if val_ref is None:
            print(f"  [{label}] série sem referência de valores — ignorada")
            continue
        values = read_range(wb, val_ref)[:n]
        values += [None] * (n - len(values))
        name = (read_cell(wb, name_ref) if name_ref else None) or cached_name or ""
        data.add_series(str(name), values, number_format="General")

    chart.replace_data(data)
    print(f"  [{label}] gráfico atualizado: {len(sers)} série(s), {n} ponto(s) "
          f"({cats[0]} a {cats[-1]})")
    return True


# ---------------------------------------------------------------------------
# Texto (células de tabela e parágrafos) preservando formatação
# ---------------------------------------------------------------------------

def set_paragraph_text(para, text: str):
    """Escreve o texto no primeiro run do parágrafo (herda a formatação
    dele) e remove os runs excedentes."""
    runs = para.runs
    if not runs:
        para.text = text
        return
    runs[0].text = text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def set_cell_text(cell, text: str):
    tf = cell.text_frame
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    set_paragraph_text(tf.paragraphs[0], text)


def set_table_data_rows(table, n: int):
    """Ajusta a tabela para ter n linhas de dados (além do cabeçalho),
    clonando a última linha ou removendo as excedentes."""
    tbl = table._tbl
    trs = tbl.tr_lst
    while len(trs) - 1 < n:
        tbl.append(deepcopy(trs[-1]))
        trs = tbl.tr_lst
    while len(trs) - 1 > n:
        tbl.remove(trs[-1])
        trs = tbl.tr_lst


# ---------------------------------------------------------------------------
# Tabelas de desempenho
# ---------------------------------------------------------------------------

def sheet_rows(wb, name):
    if name not in wb.sheetnames:
        raise KeyError(f"Aba '{name}' não encontrada na planilha")
    return [list(r) for r in wb[name].iter_rows(values_only=True)]


def update_tabela_metricas(table, wb, sheet, label):
    """Indicadores (Sharpe, Volatilidade, Beta): casa cada linha da tabela
    pelo nome do indicador na 1ª coluna."""
    dados = {str(r[0]).strip(): r[1:3] for r in sheet_rows(wb, sheet)[1:] if r[0]}
    feitos = 0
    for r in range(1, len(table.rows)):
        indicador = table.cell(r, 0).text.strip()
        if indicador not in dados:
            print(f"  [{label}] indicador '{indicador}' não achado na aba — mantido")
            continue
        eh_pct = "volatilidade" in indicador.lower()
        for c, v in enumerate(dados[indicador], start=1):
            if v is not None:
                set_cell_text(table.cell(r, c), pct(v, 2) if eh_pct else f"{v:.2f}")
        feitos += 1
    print(f"  [{label}] indicadores atualizados: {feitos} linha(s)")


def update_tabela_retornos(table, wb, sheet, label):
    """Retornos mensais: colunas = Desde o início | <ano> | Últ. 12M |
    últimos N meses (mais recente primeiro). Os rótulos dos meses vêm da
    planilha, então a janela desliza sozinha a cada mês."""
    rows = sheet_rows(wb, sheet)
    header, series = rows[1], rows[2:]
    meses_cols = [i for i, v in enumerate(header) if isinstance(v, dt.datetime)]
    n_meses_ppt = len(table.columns) - 4  # 1ª col = rótulo + 3 acumulados
    meses_cols = meses_cols[:n_meses_ppt]

    # cabeçalho: ano do YTD + rótulos dos meses
    ano = header[meses_cols[0]].year if meses_cols else dt.date.today().year
    set_cell_text(table.cell(0, 2), str(ano))
    for j, col in enumerate(meses_cols, start=4):
        set_cell_text(table.cell(0, j), mes_ano(header[col]))

    # dados: linha i da tabela <- linha i da aba (rótulos da 1ª coluna do
    # template são mantidos: "Top Ações XP" / "Ibovespa")
    for i in range(1, len(table.rows)):
        if i - 1 >= len(series):
            break
        r = series[i - 1]
        set_cell_text(table.cell(i, 1), pct(r[1]))   # desde o início
        set_cell_text(table.cell(i, 2), pct(r[2]))   # YTD
        set_cell_text(table.cell(i, 3), pct(r[3]))   # últimos 12M
        for j, col in enumerate(meses_cols, start=4):
            set_cell_text(table.cell(i, j), pct(r[col]))
    print(f"  [{label}] retornos atualizados: {len(series)} série(s), "
          f"{len(meses_cols)} mês(es) até {mes_ano(header[meses_cols[0]])}")


def update_tabela_ativos(table, wb, sheet, label):
    """Desempenho por ativo: uma linha por papel, no número que vier da
    planilha. Colunas da aba: Companhia | Ticker | Setor | Peso | Data de
    entrada | desde a entrada | no mês | no ano."""
    ativos = [r for r in sheet_rows(wb, sheet)[1:] if r and r[0]]
    set_table_data_rows(table, len(ativos))
    for i, r in enumerate(ativos, start=1):
        set_cell_text(table.cell(i, 0), str(r[0]))
        set_cell_text(table.cell(i, 1), str(r[1]))
        set_cell_text(table.cell(i, 2), pct(r[3], 2))
        set_cell_text(table.cell(i, 3), mes_ano(r[4]) if r[4] else "")
        set_cell_text(table.cell(i, 4), pct(r[5]))
        set_cell_text(table.cell(i, 5), pct(r[6]))
        set_cell_text(table.cell(i, 6), pct(r[7]))
    print(f"  [{label}] desempenho por ativo atualizado: {len(ativos)} papel(is)")


# ---------------------------------------------------------------------------
# Tabelas de composição e teses (planilha de composição, à parte)
# ---------------------------------------------------------------------------

def read_composicao(path, lang="PT"):
    """Lê a planilha de composição. Devolve (linhas_display, linhas_flat):

    * display — aba `PT`, com as células de Segmento/Setor/Peso do setor em
      branco onde há agrupamento (espelha o layout da tabela do slide 1);
    * flat — aba `PT_data`, uma linha por papel com todas as colunas
      preenchidas (usada na tabela de teses).

    Colunas (0-8): Segmento, Setor, Peso setor (Ibov), Peso setor (Carteira),
    Companhia, Ticker, Peso, Rating, Preço-Alvo.
    """
    wb = openpyxl.load_workbook(path, data_only=True)

    def dados_apos_cabecalho(rows):
        # o cabeçalho começa com "Segmento" (PT) ou "Segment" (ENG); acima
        # dele pode haver título e linhas em branco. Mantém só linhas com
        # Companhia preenchida (col 4).
        hdr = next(i for i, r in enumerate(rows)
                   if r and str(r[0]).strip() in ("Segmento", "Segment"))
        return [r for r in rows[hdr + 1:] if len(r) > 4 and r[4] is not None]

    disp = [list(r) for r in wb[lang].iter_rows(values_only=True)]
    flat = [list(r) for r in wb[f"{lang}_data"].iter_rows(values_only=True)]
    return dados_apos_cabecalho(disp), dados_apos_cabecalho(flat)


def unmerge_all(table):
    """Desfaz todas as mesclagens, devolvendo uma grade simples."""
    for r in range(len(table.rows)):
        for c in range(len(table.columns)):
            cell = table.cell(r, c)
            if cell.is_merge_origin:
                cell.split()


def merge_blank_runs(table, col, ndata):
    """Mescla verticalmente, na coluna `col`, cada sequência de uma célula
    com valor seguida de células em branco — reconstruindo o agrupamento
    exatamente como a aba PT o codifica (branco = mesclado com o de cima)."""
    r = 1
    while r <= ndata:
        start = r
        r += 1
        while r <= ndata and table.cell(r, col).text.strip() == "":
            r += 1
        if r - 1 > start:
            table.cell(start, col).merge(table.cell(r - 1, col))


def _txbody(cell):
    return cell._tc.find(qnA("txBody"))


def set_cell_paras(cell, paras):
    """Substitui os parágrafos (<a:p>) de uma célula por cópias de `paras`,
    preservando o bodyPr da célula de destino."""
    tb = _txbody(cell)
    for p in tb.findall(qnA("p")):
        tb.remove(p)
    for p in paras:
        tb.append(deepcopy(p))


def set_link(cell, url, texto="Clique aqui"):
    """Mantém o run existente (com sua formatação) e ajusta texto + link."""
    para = cell.text_frame.paragraphs[0]
    if not para.runs:
        para.text = texto
    para.runs[0].text = texto
    for extra in para.runs[1:]:
        extra._r.getparent().remove(extra._r)
    para.runs[0].hyperlink.address = url


def update_tabela_composicao(table, comp_display, label):
    """Reconstrói a tabela de composição (slide 1) a partir da aba PT.
    Colunas: Segmento | Setor | Peso setor (Ibov) | Peso setor (Carteira) |
    Companhia | Ticker | Peso | Rating | Preço-Alvo."""
    ndata = len(comp_display)
    unmerge_all(table)
    set_table_data_rows(table, ndata)
    for i, row in enumerate(comp_display, start=1):
        seg, setor, peso_i, peso_c, comp, tick, peso, rating, preco = row[:9]
        set_cell_text(table.cell(i, 0), "" if seg is None else str(seg))
        set_cell_text(table.cell(i, 1), "" if setor is None else str(setor))
        set_cell_text(table.cell(i, 2), pct_br(peso_i))
        set_cell_text(table.cell(i, 3), pct_br(peso_c))
        set_cell_text(table.cell(i, 4), "" if comp is None else str(comp))
        set_cell_text(table.cell(i, 5), "" if tick is None else str(tick))
        set_cell_text(table.cell(i, 6), pct_br(peso))
        set_cell_text(table.cell(i, 7), "" if rating is None else str(rating))
        set_cell_text(table.cell(i, 8), rs_br(preco))
    # reconstrói os agrupamentos de Segmento/Setor/Peso do setor
    for col in (0, 1, 2, 3):
        merge_blank_runs(table, col, ndata)
    print(f"  [{label}] composição atualizada: {ndata} papel(is)")


def update_tabela_teses(table, comp_flat, label):
    """Atualiza a tabela de teses (slide 4): Companhia | Ticker | Peso |
    Recomendação | Preço-Alvo | Comentários | Link para tese.

    Preenche tudo a partir da composição e monta o link por ticker. Os
    comentários (editoriais) são preservados casando pelo ticker; papéis
    novos ficam com o comentário em branco para você escrever."""
    coment = {}
    for r in range(1, len(table.rows)):
        tick = table.cell(r, 1).text.strip()
        if tick:
            coment[tick] = [deepcopy(p) for p in _txbody(table.cell(r, 5)).findall(qnA("p"))]

    set_table_data_rows(table, len(comp_flat))
    novos = []
    for i, row in enumerate(comp_flat, start=1):
        comp, tick, peso, rating, preco = row[4], row[5], row[6], row[7], row[8]
        tick = str(tick)
        set_cell_text(table.cell(i, 0), "" if comp is None else str(comp))
        set_cell_text(table.cell(i, 1), tick)
        set_cell_text(table.cell(i, 2), pct_br(peso))
        set_cell_text(table.cell(i, 3), "" if rating is None else str(rating))
        set_cell_text(table.cell(i, 4), rs_br(preco))
        if tick in coment:
            set_cell_paras(table.cell(i, 5), coment[tick])
        else:
            set_cell_text(table.cell(i, 5), "")
            novos.append(tick)
        set_link(table.cell(i, 6), LINK_TESE.format(ticker=tick))

    msg = f"  [{label}] teses atualizadas: {len(comp_flat)} papel(is), links por ticker"
    if novos:
        msg += f"; comentário a preencher: {', '.join(novos)}"
    print(msg)


# ---------------------------------------------------------------------------

def update_table(table, wb, abas, comp, label) -> bool:
    """Identifica a tabela pelo cabeçalho e delega a atualização.
    `comp` é a tupla (display, flat) da planilha de composição, ou None.
    A tabela de comentários por tese é preenchida menos a coluna de texto."""
    header = [table.cell(0, c).text.strip() for c in range(len(table.columns))]
    if header and header[0].startswith("Indicador"):
        update_tabela_metricas(table, wb, abas["metricas"], label)
    elif any(h.startswith("Desde o início") for h in header):
        update_tabela_retornos(table, wb, abas["retornos"], label)
    elif header[0].startswith("Segmento"):
        if comp is None:
            print(f"  [{label}] composição sem planilha de origem — mantida")
            return False
        update_tabela_composicao(table, comp[0], label)
    elif header[0] == "Companhia" and any(h.startswith("Link para tese") for h in header):
        if comp is None:
            print(f"  [{label}] teses sem planilha de composição — mantidas")
            return False
        update_tabela_teses(table, comp[1], label)
    elif header[0] == "Companhia" and any(h.startswith("Data de entrada") for h in header):
        update_tabela_ativos(table, wb, abas["ativos"], label)
    else:
        print(f"  [{label}] tabela editorial — mantida como está")
        return False
    return True


# ---------------------------------------------------------------------------
# Datas
# ---------------------------------------------------------------------------

RE_DATA_LONGA = re.compile(r"^\s*\d{1,2} de [a-zç]+ de \d{4}\s*$", re.I)
RE_MES_TITULO = re.compile(
    r"^\s*(" + "|".join(m.capitalize() for m in MESES_EXTENSO) + r")( de)? \d{4}\s*$")


def update_datas(prs, data: dt.date):
    longa = f"{data.day} de {MESES_EXTENSO[data.month - 1]} de {data.year}"
    titulo = f"{MESES_EXTENSO[data.month - 1].capitalize()} {data.year}"
    trocas = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                # campos de data automáticos (<a:fld type="datetime...">):
                # o PowerPoint recalcula ao abrir; atualizamos o texto em
                # cache para exports que não recalculam (PDF, prévias)
                for fld in para._p.findall(f"{{{A_NS}}}fld"):
                    if (fld.get("type", "").startswith("datetime")
                            and fld.find(f"{{{A_NS}}}t") is not None):
                        fld.find(f"{{{A_NS}}}t").text = longa
                        trocas += 1
                texto = "".join(r.text for r in para.runs)
                if RE_DATA_LONGA.match(texto):
                    set_paragraph_text(para, longa)
                    trocas += 1
                elif RE_MES_TITULO.match(texto):
                    set_paragraph_text(para, titulo)
                    trocas += 1
    print(f"  datas atualizadas para '{longa}' ({trocas} ocorrência(s))")


# ---------------------------------------------------------------------------

# Caminhos padrão na rede (sobrescreva com --template/--planilha/--saida)
DIR_BASE = r"\\xpdocs\Research\Equities\Estrategia\Carteiras\Carteiras de Ações XP"
TEMPLATE_PADRAO = rf"{DIR_BASE}\Templates\Carteira Top Ações.pptx"
PLANILHA_PADRAO = rf"{DIR_BASE}\Charts Lâmina Carteiras.xlsm"
SAIDA_PADRAO = rf"{DIR_BASE}\Lãmina Completa\Top Ações XP.pptx"
DIR_OUTPUT = rf"{DIR_BASE}\output"  # onde ficam os composicao_*.xlsx


def gerar_lamina(template=TEMPLATE_PADRAO, planilha=PLANILHA_PADRAO,
                 saida=SAIDA_PADRAO, carteira="top_acoes", data=None,
                 composicao=None):
    """Gera a lâmina. Chamável direto de um notebook/interactive window:

        from gerar_ppt import gerar_lamina
        gerar_lamina()                      # caminhos padrão da rede
        gerar_lamina(data="04/08/2026")     # data específica

    `data` aceita "DD/MM/AAAA" ou um datetime.date (padrão: hoje).
    `composicao` é o caminho do composicao_*.xlsx (padrão: pasta output/);
    se o arquivo não existir, as tabelas de composição/teses são mantidas.
    Devolve o número de gráficos/tabelas atualizados.
    """
    if data is None:
        data = dt.date.today()
    elif isinstance(data, str):
        data = dt.datetime.strptime(data, "%d/%m/%Y").date()
    abas = CARTEIRAS[carteira]
    if composicao is None:
        composicao = os.path.join(DIR_OUTPUT, abas["composicao"])

    print(f"Lendo planilha: {planilha}")
    # data_only=True devolve os valores calculados salvos pelo Excel
    wb = openpyxl.load_workbook(planilha, data_only=True, read_only=False)

    comp = None
    if os.path.exists(composicao):
        print(f"Lendo composição: {composicao}")
        comp = read_composicao(composicao)
    else:
        print(f"Composição não encontrada ({composicao}) — tabelas de "
              "composição e teses serão mantidas como no template.")

    print(f"Lendo template: {template}")
    prs = Presentation(template)

    updated = 0
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            label = f"slide {slide_idx} / {shape.name}"
            if shape.has_chart:
                updated += bool(update_chart(shape.chart, wb, label))
            elif shape.has_table:
                updated += bool(update_table(shape.table, wb, abas, comp, label))

    update_datas(prs, data)

    if updated == 0:
        raise RuntimeError("Nenhum gráfico ou tabela atualizado — verifique "
                           "se o template corresponde à planilha.")

    prs.save(saida)
    print(f"OK: {updated} objeto(s) atualizado(s) -> {saida}")
    return updated


def main():
    ap = argparse.ArgumentParser(description=__doc__.split("\n")[0])
    ap.add_argument("--template", default=TEMPLATE_PADRAO,
                    help=f"PPTX de template (padrão: {TEMPLATE_PADRAO})")
    ap.add_argument("--planilha", default=PLANILHA_PADRAO,
                    help=f"XLSM/XLSX com os gráficos (padrão: {PLANILHA_PADRAO})")
    ap.add_argument("--saida", default=SAIDA_PADRAO,
                    help=f"PPTX de saída (padrão: {SAIDA_PADRAO})")
    ap.add_argument("--carteira", choices=sorted(CARTEIRAS), default="top_acoes",
                    help="qual carteira define as abas de origem das tabelas")
    ap.add_argument("--data", default=None, metavar="DD/MM/AAAA",
                    help="data exibida nos slides (padrão: hoje)")
    ap.add_argument("--composicao", default=None,
                    help="composicao_*.xlsx (padrão: pasta output/ da carteira)")
    # parse_known_args ignora argumentos extras injetados por notebooks e
    # interactive windows (ex. --f=kernel.json do Jupyter)
    args, _ = ap.parse_known_args()

    try:
        gerar_lamina(args.template, args.planilha, args.saida,
                     args.carteira, args.data, args.composicao)
    except RuntimeError as e:
        print(str(e), file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
