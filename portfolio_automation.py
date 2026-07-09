import numpy as np
import pandas as pd

# ============================================================
# 1. CARREGAMENTO DOS DADOS
# ============================================================
portfolio_names = [
    'Carteira - TOP Ações XP',
    'Carteira - TOP DIVIDENDOS XP',
    'Carteira - TOP SMALL CAPS XP',
    'Carteira - ESG XP',
]

xpqs = pd.read_excel(r"\\xpdocs\Research\Equities\Quant\_Cross Data\xpqs-sector_classification.xlsx")[['cod_ativo', 'name', 'adjusted_GICS_sector', 'sector_xp']]


market_data = pd.read_parquet(
    r"\\xpdocs\Research\Equities\Quant\_Cross Data\economatica-market_data.parquet"
)

bdr_market_data = pd.read_csv(
    r"\\xpdocs\Research\Equities\Quant\_Cross Data\raw\bdr_market_data.csv"
)
bdr_market_data.rename(columns={'Ativo': 'cod_ativo', 'Data': 'data'}, inplace=True)
bdr_market_data['cod_ativo'] = bdr_market_data['cod_ativo'].str.replace("<XBSP>", "")
bdr_market_data['data'] = pd.to_datetime(bdr_market_data['data'])

valid_dates = market_data['data'].unique()
bdr_market_data = bdr_market_data[bdr_market_data['data'].isin(valid_dates)].copy()
bdr_market_data['adj_close_price'] = pd.to_numeric(bdr_market_data['adj_close_price'], errors='coerce')
market_data = pd.concat([market_data, bdr_market_data], ignore_index=True)

indices = pd.read_parquet(r"\\xpdocs\Research\Equities\Quant\_Cross Data\economatica-indices.parquet")
indices = indices.loc[indices['cod_ativo'].isin(['Ibovespa', 'SMLL', 'ISEE'])].copy()
indices['data'] = pd.to_datetime(indices['data'])

file_path = r"\\xpdocs\Research\Equities\Estrategia\Carteiras\Performance carteiras.xlsm"

composition_dict = {}

for portfolio in portfolio_names:
    df = pd.read_excel(file_path, sheet_name=portfolio, skiprows=5)

    if portfolio == 'Carteira - TOP Ações XP':
        df = df.iloc[:, 89:]
        df.rename(columns={'Ticker.1': 'cod_ativo'}, inplace=True)
    else:
        df = df.iloc[:, 3:]
        df.rename(columns={'Ticker': 'cod_ativo'}, inplace=True)

    composition_dict[portfolio] = df


# ============================================================
# 2. PREPARAÇÃO DO MARKET DATA
# ============================================================
md = market_data.copy()
md['data'] = pd.to_datetime(md['data'])
md = md[['cod_ativo', 'data', 'adj_close_price']].dropna()
md = md.sort_values(['cod_ativo', 'data'])

# retorno diário por ativo
md['ret'] = md.groupby('cod_ativo')['adj_close_price'].pct_change()

# vetor de datas úteis que existem no market_data
datas_md = np.sort(md['data'].unique())


# ============================================================
# 3. FUNÇÃO DE PERFORMANCE — REBALANCEAMENTO MENSAL (BASE 100)
#    HÍBRIDA: sem drift até maio/2026, com drift a partir de jun/2026
# ============================================================
DRIFT_INICIO = pd.Timestamp(2026, 5, 29)   # rebals >= esta data usam drift


def _ret_diario_sem_drift(ret_periodo, comp_periodo):
    """Pesos renormalizados ao alvo todo dia (metodologia antiga)."""
    rp = ret_periodo.merge(comp_periodo, on='cod_ativo', how='inner')

    def retorno_diario(grupo):
        w = grupo['peso'] / grupo['peso'].sum()
        return (w * grupo['ret']).sum()

    return rp.groupby('data').apply(retorno_diario)


def _ret_diario_com_drift(ret_periodo, comp_periodo):
    """
    Buy & hold dentro da janela: o peso evolui com os preços.
    Para cada dia t:
        ret_carteira_t = Σ_i ( w_i,t * ret_i,t )
    onde w_i,t é o peso DRIFTED no início do dia (soma 1), atualizado por:
        w_i,t = w_i,t-1 * (1 + ret_i,t)  -> renormaliza só para obter a fração do dia
    Ativo sem retorno no dia (NaN) é tratado como ret 0 (carrega o peso).
    """
    # pesos-alvo iniciais (normalizados), indexados por ativo
    pesos = comp_periodo.set_index('cod_ativo')['peso']
    pesos = pesos / pesos.sum()

    # matriz dias x ativos com os retornos da janela
    rp = ret_periodo[ret_periodo['cod_ativo'].isin(pesos.index)].copy()
    ret_mat = (rp.pivot_table(index='data', columns='cod_ativo', values='ret')
                 .reindex(columns=pesos.index)
                 .sort_index())

    datas = ret_mat.index
    out = pd.Series(index=datas, dtype=float)

    # capital relativo de cada ativo; começa nos pesos-alvo
    cap = pesos.copy()

    for dia in datas:
        r = ret_mat.loc[dia].fillna(0.0)      # ret do dia (NaN -> 0)
        w = cap / cap.sum()                   # peso efetivo no início do dia (soma 1)
        out.loc[dia] = float((w * r).sum())   # retorno da carteira no dia
        cap = cap * (1.0 + r)                 # deixa o capital derivar (sem renormalizar)

    return out


def calcular_performance(composition):
    comp = composition.copy()

    comp_long = comp.melt(id_vars='cod_ativo', var_name='data_rebal', value_name='peso')
    comp_long['data_rebal'] = pd.to_datetime(comp_long['data_rebal'])
    comp_long = comp_long.dropna(subset=['peso'])

    rebal_dates = sorted(comp_long['data_rebal'].unique())
    retornos_carteira = []

    for i, start in enumerate(rebal_dates):
        next_start = (rebal_dates[i + 1] if i + 1 < len(rebal_dates)
                      else md['data'].max())

        comp_periodo = comp_long[comp_long['data_rebal'] == start][['cod_ativo', 'peso']]

        mask = (md['data'] > start) & (md['data'] <= next_start)
        ret_periodo = md.loc[mask, ['cod_ativo', 'data', 'ret']]

        # escolhe a metodologia conforme a data do rebal
        if pd.Timestamp(start) >= DRIFT_INICIO:
            ret_dia = _ret_diario_com_drift(ret_periodo, comp_periodo)
        else:
            ret_dia = _ret_diario_sem_drift(ret_periodo, comp_periodo)

        retornos_carteira.append(ret_dia)

    serie_ret = pd.concat(retornos_carteira).sort_index()
    serie_ret = serie_ret[~serie_ret.index.duplicated(keep='first')]

    fatores = (1 + serie_ret.fillna(0)).cumprod()

    data_inicial = pd.Timestamp(rebal_dates[0])

    base100 = pd.concat([
        pd.Series([100.0], index=[data_inicial]),
        100 * fatores
    ]).sort_index()

    return base100


# ============================================================
# 3b. FUNÇÃO AUXILIAR — ÍNDICE EM BASE 100, ALINHADO À CARTEIRA
# ============================================================
def indice_base100(cod_indice, datas_carteira):
    """
    Retorna o índice em base 100, alinhado exatamente às datas da carteira
    (reindex + ffill) e ancorado em 100 no inception da carteira.
    """
    serie = (indices.loc[indices['cod_ativo'] == cod_indice, ['data', 'close_price']]
                    .drop_duplicates('data')
                    .set_index('data')['close_price']
                    .sort_index())

    inception = datas_carteira.min()
    fim = datas_carteira.max()

    serie = serie[(serie.index >= inception - pd.Timedelta(days=10)) &
                  (serie.index <= fim)]

    if serie.empty:
        return pd.Series(dtype=float)

    serie = serie.reindex(datas_carteira.union(serie.index)).ffill()
    serie = serie.reindex(datas_carteira)

    primeiro_valido = serie.first_valid_index()
    if primeiro_valido is None:
        return pd.Series(dtype=float)

    return 100 * serie / serie.loc[primeiro_valido]


# ============================================================
# 4. APLICA NAS CARTEIRAS E COMPARA COM BENCHMARKS
# ============================================================
benchmarks_por_carteira = {
    'Carteira - TOP Ações XP':        ['Ibovespa'],
    'Carteira - TOP DIVIDENDOS XP':   ['Ibovespa'],
    'Carteira - TOP SMALL CAPS XP':   ['Ibovespa', 'SMLL'],   
    'Carteira - ESG XP':              ['Ibovespa', 'ISEE'],
}

resultado_dfs = {}

for portfolio in portfolio_names:
    curva_carteira = calcular_performance(composition_dict[portfolio])

    df_port = pd.DataFrame({portfolio: curva_carteira})

    for bench in benchmarks_por_carteira[portfolio]:
        df_port[bench] = indice_base100(bench, curva_carteira.index)

    df_port.index.name = 'data'
    resultado_dfs[portfolio] = df_port


df_top_acoes      = resultado_dfs['Carteira - TOP Ações XP']
df_top_dividendos = resultado_dfs['Carteira - TOP DIVIDENDOS XP']
df_top_smallcaps  = resultado_dfs['Carteira - TOP SMALL CAPS XP']
df_esg            = resultado_dfs['Carteira - ESG XP']

# ============================================================
# 4b. EXPORTA AS CURVAS BASE 100 (input para os scripts de e-mail/métricas)
# ============================================================
output_dir = r"\\xpdocs\Research\Equities\Estrategia\Carteiras\Carteiras de Ações XP\output"
 
# nome de arquivo por carteira (segue o padrão pedido: <nome>_base_100.xlsx)
arq_base100 = {
    'Carteira - TOP Ações XP':      'top_acoes_base_100',
    'Carteira - TOP DIVIDENDOS XP': 'top_dividendos_base_100',
    'Carteira - TOP SMALL CAPS XP': 'top_small_caps_base_100',
    'Carteira - ESG XP':            'esg_base_100',
}
 
for portfolio in portfolio_names:
    df_b100 = resultado_dfs.get(portfolio)
    if df_b100 is None:
        print(f"[AVISO] base 100 não exportada (carteira não processada): {portfolio}")
        continue
 
    caminho = f"{output_dir}\\{arq_base100[portfolio]}.xlsx"
 
    with pd.ExcelWriter(caminho, engine='xlsxwriter') as writer:
        # exporta a curva base 100 (carteira + benchmarks) com a data no índice
        df_b100.to_excel(writer, sheet_name='base100')
 
        wb = writer.book
        ws = writer.sheets['base100']
 
        fmt_data  = wb.add_format({'num_format': 'dd/mm/yyyy', 'align': 'center'})
        fmt_num   = wb.add_format({'num_format': '0.00', 'align': 'center'})
 
        # coluna A = datas (índice)
        ws.set_column(0, 0, 12, fmt_data)
        # demais colunas = valores base 100
        ws.set_column(1, len(df_b100.columns), 16, fmt_num)
 
    print(f"Base 100 exportada: {caminho}")


# ============================================================
# 4c. SHARPE / VOLATILIDADE / BETA (últimos 12 meses)
#     Carteira vs. seus benchmarks | 1 arquivo por carteira
# ============================================================

indices = pd.read_parquet(r"\\xpdocs\Research\Equities\Quant\_Cross Data\economatica-indices.parquet")

# CDI (Sharpe)
cdi = (indices.loc[indices['cod_ativo'] == 'CDI Acumulado', ['data', 'close_price']]
              .drop_duplicates('data')
              .set_index('data')['close_price']
              .sort_index())


def _ret_diarios(serie):
    return serie.dropna().pct_change().dropna()


def _ultimos_12m(serie):
    """Recorta a série aos últimos 12 meses (a partir da última data disponível)."""
    s = serie.dropna()
    if s.empty:
        return s
    corte = s.index.max() - pd.DateOffset(months=12)
    return s[s.index >= corte]


def indicadores_12m(df_port, serie_cdi):
    """Sharpe, Volatilidade e Beta (últimos 12m) p/ carteira e cada benchmark. Beta vs. Ibovespa."""
    ret_ibov = _ret_diarios(_ultimos_12m(df_port['Ibovespa'])) if 'Ibovespa' in df_port.columns else None
    ret_cdi  = _ret_diarios(_ultimos_12m(serie_cdi))

    def _vol(r):
        return r.std() * np.sqrt(252)

    def _sharpe(r):
        idx = r.index.intersection(ret_cdi.index)
        if len(idx) < 2:
            return np.nan
        v = r.loc[idx].std() * np.sqrt(252)
        return ((r.loc[idx] - ret_cdi.loc[idx]).mean() * 252) / v if v > 0 else np.nan

    def _beta(r, eh_ibov):
        if eh_ibov:
            return 1.0                      
        if ret_ibov is None:
            return np.nan
        idx = r.index.intersection(ret_ibov.index)
        rc, ri = r.loc[idx], ret_ibov.loc[idx]
        var_i = np.var(ri)
        return np.cov(rc, ri)[0, 1] / var_i if var_i > 0 else np.nan

    dados = {col: [_sharpe(r := _ret_diarios(_ultimos_12m(df_port[col]))),
                   _vol(r),
                   _beta(r, col == 'Ibovespa')]
             for col in df_port.columns}

    return pd.DataFrame(dados, index=['Sharpe', 'Volatilidade', 'Beta'])


arq_fig2 = {
    'Carteira - TOP Ações XP':      'portfolio_metrics_top_acoes',
    'Carteira - TOP DIVIDENDOS XP': 'portfolio_metrics_top_dividendos',
    'Carteira - TOP SMALL CAPS XP': 'portfolio_metrics_top_small_caps',
    'Carteira - ESG XP':            'portfolio_metrics_esg',
}

for portfolio in portfolio_names:
    tab = indicadores_12m(resultado_dfs[portfolio], cdi)
    caminho = f"{output_dir}\\{arq_fig2[portfolio]}.xlsx"
    tab.to_excel(caminho)
    print(f"Exportado: {caminho}")

# ============================================================
# 4d. VERSÃO "LÂMINA" DOS RESULTADOS
#     Na Small Caps, a lâmina/PPT usa APENAS o SMLL como benchmark
#     (o Ibovespa entra em todo o resto do fluxo, menos na lâmina).
# ============================================================
BENCH_LAMINA = {
    'Carteira - TOP SMALL CAPS XP': ['SMLL'],   # remove Ibovespa da lâmina/PPT
}

def _df_para_lamina(portfolio):
    """Retorna o df_port para uso na lâmina/PPT, aplicando o filtro de
    benchmarks específico da lâmina quando houver (senão, usa o padrão)."""
    df = resultado_dfs[portfolio]
    benchs_lamina = BENCH_LAMINA.get(portfolio)
    if benchs_lamina is None:
        return df
    nome_cart = _nome_col_carteira(df)
    cols_manter = [nome_cart] + [b for b in benchs_lamina if b in df.columns]
    return df[cols_manter]

# ============================================================
# 5. TABELA DE RETORNOS: Since Inception, YTD, LTM, 12 meses, anuais
# ============================================================
def _ret_periodo(serie, data_ini, data_fim):
    """Retorno simples entre dois pontos da série base 100 (usa ffill)."""
    s = serie.dropna()
    if s.empty:
        return np.nan

    s_fim = s[s.index <= data_fim]
    if s_fim.empty:
        return np.nan
    v_fim = s_fim.iloc[-1]

    s_ini = s[s.index <= data_ini]
    if s_ini.empty:
        v_ini = s.iloc[0]
    else:
        v_ini = s_ini.iloc[-1]

    return v_fim / v_ini - 1


def _ret_mes(serie, ano, mes):
    """Retorno do mês cheio (fim do mês anterior -> fim do mês)."""
    s = serie.dropna()
    if s.empty:
        return np.nan

    ini = pd.Timestamp(ano, mes, 1) - pd.Timedelta(days=1)
    fim = pd.Timestamp(ano, mes, 1) + pd.offsets.MonthEnd(0)

    v_ini = s[s.index <= ini]
    v_fim = s[s.index <= fim]
    if v_ini.empty or v_fim.empty:
        return np.nan
    return v_fim.iloc[-1] / v_ini.iloc[-1] - 1


def _ret_ano(serie, ano):
    """
    Retorno do ano cheio (fim do ano anterior -> fim do ano).
    No ano de inception, usa o inception como ponto inicial.
    """
    s = serie.dropna()
    if s.empty:
        return np.nan

    inception = s.index.min()
    ini = pd.Timestamp(ano, 1, 1) - pd.Timedelta(days=1)
    fim = pd.Timestamp(ano, 12, 31)

    if ini < inception:
        v_ini = s.iloc[0]
    else:
        v = s[s.index <= ini]
        if v.empty:
            return np.nan
        v_ini = v.iloc[-1]

    v_fim = s[s.index <= fim]
    if v_fim.empty:
        return np.nan
    return v_fim.iloc[-1] / v_ini - 1


def tabela_retornos(df_port, n_meses=12):
    """
    Gera a tabela de retornos para todas as colunas (carteira + benchmarks):
    Since inception, YTD, LTM, últimos N meses, [coluna vazia], retornos anuais.
    """
    data_fim = df_port.dropna(how='all').index.max()
    ano_atual = data_fim.year

    inception   = df_port.index.min()
    inicio_ytd  = pd.Timestamp(ano_atual, 1, 1) - pd.Timedelta(days=1)
    inicio_ltm  = data_fim - pd.DateOffset(years=1)

    meses = []
    ref = pd.Timestamp(data_fim.year, data_fim.month, 1)
    for i in range(n_meses):
        m = ref - pd.DateOffset(months=i)
        meses.append((m.year, m.month))

    ano_inception = inception.year
    anos = list(range(ano_atual - 1, ano_inception - 1, -1))

    linhas = {}
    for col in df_port.columns:
        s = df_port[col]
        row = {
            'Since inception': _ret_periodo(s, inception, data_fim),
            'YTD':             _ret_periodo(s, inicio_ytd, data_fim),
            'LTM':             _ret_periodo(s, inicio_ltm, data_fim),
        }
        for (ano, mes) in meses:
            label = pd.Timestamp(ano, mes, 1).strftime('%b-%y')
            row[label] = _ret_mes(s, ano, mes)

        row[''] = np.nan

        for ano in anos:
            row[str(ano)] = _ret_ano(s, ano)

        linhas[col] = row

    tabela = pd.DataFrame(linhas).T

    col_meses = [pd.Timestamp(a, m, 1).strftime('%b-%y') for (a, m) in meses]
    col_anos  = [str(a) for a in anos]
    tabela = tabela[['Since inception', 'YTD', 'LTM'] + col_meses + [''] + col_anos]

    tabela_fmt = tabela.copy()
    for c in tabela_fmt.columns:
        if c == '':
            tabela_fmt[c] = ''
        else:
            tabela_fmt[c] = (tabela[c] * 100).round(1).astype(str) + '%'
            tabela_fmt[c] = tabela_fmt[c].replace('nan%', '')

    return tabela, tabela_fmt


# ============================================================
# 6. GERA AS TABELAS DAS CARTEIRAS (export numérico, % formatado no Excel)
# ============================================================
tab_acoes_raw,      tab_acoes      = tabela_retornos(df_top_acoes)
tab_dividendos_raw, tab_dividendos = tabela_retornos(df_top_dividendos)
tab_smallcaps_raw,  tab_smallcaps  = tabela_retornos(df_top_smallcaps)
tab_esg_raw,        tab_esg        = tabela_retornos(df_esg)

tab_acoes_raw      = tab_acoes_raw.rename(index={'Carteira - TOP Ações XP': 'Top Ideas'})
tab_dividendos_raw = tab_dividendos_raw.rename(index={'Carteira - TOP DIVIDENDOS XP': 'Top Dividends'})
tab_smallcaps_raw  = tab_smallcaps_raw.rename(index={'Carteira - TOP SMALL CAPS XP': 'Top Small Caps'})
tab_esg_raw        = tab_esg_raw.rename(index={'Carteira - ESG XP': 'Top ESG'})

output_dir = r"\\xpdocs\Research\Equities\Estrategia\Carteiras\Carteiras de Ações XP\output"


import datetime

def _col_para_data(col):
    """
    Se a coluna for um rótulo de mês no formato '%b-%y' (ex: 'Jun-25'),
    retorna um datetime (1º dia do mês: 01/Jun/25). Senão, retorna None.
    """
    if not isinstance(col, str) or col == '':
        return None
    try:
        return datetime.datetime.strptime(col, '%b-%y')
    except ValueError:
        return None


def exportar_tabela_retornos(tab_raw, caminho, sheet='performance'):
    """
    Exporta a tabela de retornos:
      - Header dos meses gravado como DATA real (1º dia do mês), formato 'mmm/yy'.
      - Colunas de texto (Since inception, YTD, LTM, anos, separador '') ficam texto.
      - Valores em % (número de verdade, formato 0.0%).
    """
    tab = tab_raw.copy()

    with pd.ExcelWriter(caminho, engine='xlsxwriter') as writer:
        tab.to_excel(writer, sheet_name=sheet, startrow=1, header=False)
        wb = writer.book
        ws = writer.sheets[sheet]

        fmt_pct  = wb.add_format({'num_format': '0.0%', 'align': 'center'})
        fmt_date = wb.add_format({'num_format': 'mmm/yy', 'align': 'center', 'bold': True})
        fmt_txt  = wb.add_format({'align': 'center', 'bold': True})

        ws.write(0, 0, '')

        for j, col in enumerate(tab.columns, start=1):
            dt = _col_para_data(col)
            if dt is not None:
                ws.write_datetime(0, j, dt, fmt_date)
            else:
                ws.write(0, j, str(col), fmt_txt)

        for j, col in enumerate(tab.columns, start=1):
            if col == '':
                continue
            ws.set_column(j, j, 9, fmt_pct)

        ws.set_column(0, 0, 16)

    print(f"Performance exportada: {caminho}")


exportar_tabela_retornos(tab_acoes_raw,      f"{output_dir}\\tab_performance_top_acoes.xlsx")
exportar_tabela_retornos(tab_dividendos_raw, f"{output_dir}\\tab_performance_top_dividendos.xlsx")
exportar_tabela_retornos(tab_smallcaps_raw,  f"{output_dir}\\tab_performance_top_small_caps.xlsx")
exportar_tabela_retornos(tab_esg_raw,        f"{output_dir}\\tab_performance_esg.xlsx")


# ============================================================
# 7. TABELA POR COMPONENTE — atual (mês aberto) e último rebal (mês fechado)
# ============================================================
mapa_info = xpqs.drop_duplicates('cod_ativo').set_index('cod_ativo')
mapa_nome   = mapa_info['name'].to_dict()
mapa_setor  = mapa_info['sector_xp'].to_dict()

def _primeiro_dia_util_do_mes(data, datas_pregao):
    """
    Mês de entrada para EXIBIÇÃO, considerando TRADING DAYS reais (com feriados):
      - se data_rebal == último TRADING DAY do seu mês  -> entra no mês seguinte
      - se for no meio do mês                           -> entra no próprio mês
    Retorna o 1º dia (1º do mês de entrada).

    `datas_pregao`: array/Series de datas reais de pregão (ex: md['data'].unique()).
    """
    d = pd.to_datetime(data)
    dp = pd.DatetimeIndex(sorted(pd.to_datetime(datas_pregao)))

    # último TRADING DAY real do mês de d (maior pregão <= fim do mês)
    fim_mes = pd.Timestamp(d.year, d.month, 1) + pd.offsets.MonthEnd(0)
    pregoes_do_mes = dp[(dp >= pd.Timestamp(d.year, d.month, 1)) & (dp <= fim_mes)]
    if len(pregoes_do_mes) == 0:
        ultimo_trading = fim_mes   # fallback
    else:
        ultimo_trading = pregoes_do_mes[-1]

    # se rebal == último trading day -> mês seguinte; senão -> mesmo mês
    if d.normalize() == ultimo_trading.normalize():
        return pd.Timestamp(d.year, d.month, 1) + pd.offsets.MonthBegin(1)
    else:
        return pd.Timestamp(d.year, d.month, 1)

def _preco_em(cod_ativo, data_alvo):
    """Último preço ajustado disponível <= data_alvo para o ativo."""
    s = md[(md['cod_ativo'] == cod_ativo) & (md['data'] <= data_alvo)]
    if s.empty:
        return np.nan
    return s.sort_values('data')['adj_close_price'].iloc[-1]


def _rotulo_mes_entrada(data_entrada):
    """
    Rótulo da data de entrada (só exibição, não altera cálculo).
    Se a entrada cair no ÚLTIMO DIA ÚTIL do mês, exibe o mês SUBSEQUENTE.
    """
    data = pd.to_datetime(data_entrada)

    ultimo_dia_util = data + pd.offsets.MonthEnd(0)
    if ultimo_dia_util.weekday() >= 5:
        ultimo_dia_util -= pd.offsets.BDay(1)

    if data.normalize() == ultimo_dia_util.normalize():
        return (data + pd.offsets.MonthBegin(1)).strftime('%b-%y')
    return data.strftime('%b-%y')


def _data_entrada_continua(comp_long, cod_ativo, rebal_dates, rebal_ref):
    """
    Data de entrada na carteira considerando o período contínuo que inclui rebal_ref.
    """
    presentes = set(comp_long.loc[comp_long['cod_ativo'] == cod_ativo, 'data_rebal'])
    idx_ref = rebal_dates.index(rebal_ref)

    entrada = rebal_ref
    for i in range(idx_ref, -1, -1):
        if rebal_dates[i] in presentes:
            entrada = rebal_dates[i]
        else:
            break
    return entrada


def tabela_componentes(composition, rebal_ref, ini_mes, fim_mes, ini_ytd, fim_ref):
    """Monta a tabela por componente para o rebalanceamento `rebal_ref`."""
    comp = composition.copy()
    comp_long = comp.melt(id_vars='cod_ativo', var_name='data_rebal', value_name='peso')
    comp_long['data_rebal'] = pd.to_datetime(comp_long['data_rebal'])
    comp_long = comp_long.dropna(subset=['peso'])

    rebal_dates = sorted(comp_long['data_rebal'].unique())

    comp_ref = comp_long[comp_long['data_rebal'] == rebal_ref][['cod_ativo', 'peso']].copy()

    linhas = []
    for _, r in comp_ref.iterrows():
        cod = r['cod_ativo']
        peso = r['peso'] / comp_ref['peso'].sum()

        entrada = _data_entrada_continua(comp_long, cod, rebal_dates, rebal_ref)

        p_entrada = _preco_em(cod, entrada)
        p_fim_ref = _preco_em(cod, fim_ref)
        desemp_entrada = (p_fim_ref / p_entrada - 1) if (p_entrada and p_fim_ref) else np.nan

        p_ini_mes = _preco_em(cod, ini_mes)
        p_fim_mes = _preco_em(cod, fim_mes)
        desemp_mes = (p_fim_mes / p_ini_mes - 1) if (p_ini_mes and p_fim_mes) else np.nan

        p_ini_ytd = _preco_em(cod, ini_ytd)
        p_fim_ytd = _preco_em(cod, fim_ref)
        desemp_ytd = (p_fim_ytd / p_ini_ytd - 1) if (p_ini_ytd and p_fim_ytd) else np.nan
        
        linhas.append({
            'Companhia':                 mapa_nome.get(cod, ''),
            'Ticker':                    cod,
            'Setor':                     mapa_setor.get(cod, ''),
            'Peso':                      peso,
            'Data de entrada':           _primeiro_dia_util_do_mes(entrada, md['data'].unique()),  # data real (1º dia útil), exibida como mmm/yy
            'Desempenho desde entrada':  desemp_entrada,   # cálculo usa 'entrada' (último trading day)
            'Desempenho no mês':         desemp_mes,
            'Desempenho YTD':            desemp_ytd,
        })

    tabela = pd.DataFrame(linhas).sort_values('Ticker', ascending=True).reset_index(drop=True)
    return tabela


def _formata_componentes(tabela):
    """Formata pesos e desempenhos em %."""
    t = tabela.copy()
    t['Peso'] = (t['Peso'] * 100).round(1).astype(str) + '%'
    for c in ['Desempenho desde entrada', 'Desempenho no mês', 'Desempenho YTD']:
        t[c] = (tabela[c] * 100).round(1).astype(str) + '%'
        t[c] = t[c].replace('nan%', '')
    return t


def gerar_tabelas_componentes(composition):
    """Gera as duas tabelas (atual e último rebal) para uma carteira."""
    comp = composition.copy()
    comp_long = comp.melt(id_vars='cod_ativo', var_name='data_rebal', value_name='peso')
    comp_long['data_rebal'] = pd.to_datetime(comp_long['data_rebal'])
    comp_long = comp_long.dropna(subset=['peso'])
    rebal_dates = sorted(comp_long['data_rebal'].unique())

    rebal_atual = rebal_dates[-1]
    rebal_anterior = rebal_dates[-2] if len(rebal_dates) >= 2 else rebal_dates[-1]

    data_max = md['data'].max()

    ini_mes_atual = pd.Timestamp(data_max.year, data_max.month, 1) - pd.Timedelta(days=1)
    fim_mes_atual = data_max
    ini_ytd_atual = pd.Timestamp(data_max.year, 1, 1) - pd.Timedelta(days=1)

    tab_atual = tabela_componentes(
        comp, rebal_atual,
        ini_mes=ini_mes_atual, fim_mes=fim_mes_atual,
        ini_ytd=ini_ytd_atual, fim_ref=data_max
    )

    mes_ant = pd.Timestamp(data_max.year, data_max.month, 1) - pd.Timedelta(days=1)
    ini_mes_ant = pd.Timestamp(mes_ant.year, mes_ant.month, 1) - pd.Timedelta(days=1)
    fim_mes_ant = mes_ant
    ini_ytd_ant = pd.Timestamp(mes_ant.year, 1, 1) - pd.Timedelta(days=1)

    tab_ultimo = tabela_componentes(
        comp, rebal_anterior,
        ini_mes=ini_mes_ant, fim_mes=fim_mes_ant,
        ini_ytd=ini_ytd_ant, fim_ref=fim_mes_ant
    )

    return tab_atual, tab_ultimo


# ============================================================
# 8. GERA E EXPORTA AS TABELAS POR COMPONENTE (números de verdade)
# ============================================================
output_dir = r"\\xpdocs\Research\Equities\Estrategia\Carteiras\Carteiras de Ações XP\output"

mapa_arquivo = {
    'Carteira - TOP Ações XP':      'top_acoes',
    'Carteira - TOP DIVIDENDOS XP': 'top_dividendos',
    'Carteira - TOP SMALL CAPS XP': 'top_small_caps',
    'Carteira - ESG XP':            'esg',
}

# colunas que são percentuais
cols_pct = ['Peso', 'Desempenho desde entrada', 'Desempenho no mês', 'Desempenho YTD']

tabelas_componentes = {}

for portfolio in portfolio_names:
    tab_atual, tab_ultimo = gerar_tabelas_componentes(composition_dict[portfolio])

    tabelas_componentes[portfolio] = {'atual': tab_atual, 'ultimo': tab_ultimo}

    nome = mapa_arquivo[portfolio]

    for sufixo, tab in [('atual', tab_atual), ('ultimo_rebal', tab_ultimo)]:
        caminho = f"{output_dir}\\componentes_{nome}_{sufixo}.xlsx"

        with pd.ExcelWriter(caminho, engine='xlsxwriter') as writer:
            tab.to_excel(writer, index=False, sheet_name='componentes')
            wb = writer.book
            ws = writer.sheets['componentes']

            fmt_pct = wb.add_format({'num_format': '0.0%'})

            for c in cols_pct:
                col_idx = tab.columns.get_loc(c)
                ws.set_column(col_idx, col_idx, 14, fmt_pct)


# ============================================================
# 9. TABELA DE COMPOSIÇÃO (ENG + PT) por carteira
#    - Composição: Performance carteiras.xlsm (último rebal)
#    - Pesos setoriais do Ibovespa: calculados na hora (sem Excel pronto)
#    - Segmento: derivado do sector_dict
#    - Abas VISUAIS (com merge) + abas FLAT (sem merge, para query)
# ============================================================
from pathlib import Path
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, Border, Side, PatternFill

output_dir = r"\\xpdocs\Research\Equities\Estrategia\Carteiras\Carteiras de Ações XP\output"

# ------------------------------------------------------------
# 9.0  Mapas de setor/segmento e traduções
# ------------------------------------------------------------
sector_dict = {
    'Commodities':          ['Energy', 'Materials'],
    'Domestic Defensives':  ['Consumer Staples', 'Health Care', 'Utilities'],
    'Domestic Cyclicals':   ['Consumer Discretionary', 'Industrials',
                             'Information Technology', 'Real Estate', 'Communication Services'],
    'Financials':           ['Financials'],
}
sector_to_segment = {s: grupo for grupo, setores in sector_dict.items() for s in setores}

SETOR_PT = {
    'Energy': 'Energia',
    'Materials': 'Materiais',
    'Industrials': 'Bens Industriais',
    'Consumer Discretionary': 'Consumo Discricionário',
    'Consumer Staples': 'Consumo Básico',
    'Health Care': 'Saúde',
    'Financials': 'Financeiro',
    'Information Technology': 'Tecnologia da Informação',
    'Communication Services': 'Comunicações',
    'Utilities': 'Utilidade Pública',
    'Real Estate': 'Imobiliário',
}
SEGMENTO_PT = {
    'Commodities': 'Commodities',
    'Domestic Defensives': 'Defensivas Domésticas',
    'Domestic Cyclicals': 'Cíclicas Domésticas',
    'Financials': 'Financeiro',
}
ABREV_EN = {
    'Consumer Discretionary': 'Cons. Disc.',
    'Information Technology': 'Info Tech.',
    'Consumer Staples': 'Cons. Staples',
}
ABREV_PT = {
    'Consumo Discricionário': 'Cons. Disc.',
    'Tecnologia da Informação': 'Tec. Info.',
    'Consumo Básico': 'Cons. Básico',
}

# ------------------------------------------------------------
# 9.1  Bases auxiliares (metadados de cada ativo)
# ------------------------------------------------------------
comp_sheet = pd.read_excel(
    r'\\xpdocs\Research\Equities\COMP SHEET\raw_data.xlsx', sheet_name='Sheet1'
)[['TICKER', 'NAME', 'TARGET', 'RECOMMENDATION']].rename(columns={
    'TICKER': 'Ticker', 'NAME': 'Company', 'TARGET': 'Target Price', 'RECOMMENDATION': 'Rating'
})

sector_classification_raw = xpqs[['cod_ativo', 'adjusted_GICS_sector']].rename(
    columns={'cod_ativo': 'Ticker', 'adjusted_GICS_sector': 'Sector'}
).drop_duplicates('Ticker')

# ------------------------------------------------------------
# 9.2  Pesos setoriais do Ibovespa (calculados na hora)
# ------------------------------------------------------------
def _download_ibov_composition():
    fp = Path(r"\\xpdocs\research\equities\Quant\_Cross Data\economatica-index_composition.csv")
    df = (pd.read_csv(fp, usecols=['cod_ativo', 'data', 'IBOV'], sep=';')
            .dropna()
            .rename(columns={'IBOV': 'weight'}))
    df['weight'] = df['weight'].str.replace(',', '.').astype(float)
    df.sort_values(['data', 'cod_ativo'], inplace=True)
    return df


def calcular_pesos_ibovespa():
    """
    Peso setorial (GICS individual) + por grupo do Ibovespa na data mais recente.
    Retorna (individual, grupos), ambos com pesos em fração (0-1).
    """
    members = _download_ibov_composition()

    setores = (xpqs[['cod_ativo', 'adjusted_GICS_sector']]
               .drop_duplicates('cod_ativo')
               .set_index('cod_ativo'))

    df = members.merge(setores, on='cod_ativo', how='left')
    df['sector_group'] = df['adjusted_GICS_sector'].map(sector_to_segment)

    ultima = df['data'].max()
    df_u = df[df['data'] == ultima]

    individual = (df_u.groupby('adjusted_GICS_sector')['weight'].sum() / 100).reset_index()
    individual = individual.rename(columns={
        'adjusted_GICS_sector': 'Sector', 'weight': 'Sector Weight (Ibovespa)'
    })

    grupos = (df_u.groupby('sector_group')['weight'].sum() / 100).reset_index()
    grupos = grupos.rename(columns={
        'sector_group': 'Segment', 'weight': 'Segment Weight (Ibovespa)'
    })

    return individual, grupos


sector_weight_ibovespa, segment_weight_ibovespa = calcular_pesos_ibovespa()

segment_sector = pd.DataFrame(
    [{'Sector': s, 'Segment': seg} for s, seg in sector_to_segment.items()]
)

with pd.ExcelWriter(f"{output_dir}\\sector_weights_ibovespa.xlsx") as writer:
    sector_weight_ibovespa.to_excel(writer, sheet_name='Individual Sectors', index=False)
    segment_weight_ibovespa.to_excel(writer, sheet_name='Sector Groups', index=False)

# ------------------------------------------------------------
# 9.3  Monta o DataFrame de composição de uma carteira
# ------------------------------------------------------------
def _composicao_vigente(composition):
    """Ticker + Weight do último rebalanceamento (carteira vigente)."""
    comp_long = composition.melt(id_vars='cod_ativo', var_name='data_rebal', value_name='peso')
    comp_long['data_rebal'] = pd.to_datetime(comp_long['data_rebal'])
    comp_long = comp_long.dropna(subset=['peso'])
    ult = comp_long['data_rebal'].max()
    comp = comp_long[comp_long['data_rebal'] == ult][['cod_ativo', 'peso']].copy()
    comp = comp.rename(columns={'cod_ativo': 'Ticker', 'peso': 'Weight'})
    comp['Weight'] = comp['Weight'] / comp['Weight'].sum()
    return comp


def montar_df_composicao(composition, idioma='EN'):
    df = _composicao_vigente(composition)
    df = df.merge(comp_sheet, on='Ticker', how='left')
    df = df.merge(sector_classification_raw, on='Ticker', how='left')   # setor GICS puro
    df = df.merge(sector_weight_ibovespa, on='Sector', how='left')      # casa em GICS puro
    df = df.merge(segment_sector, on='Sector', how='left')              # casa em GICS puro
    df['Sector Weight (Portfolio)'] = df.groupby('Sector')['Weight'].transform('sum')

    df = df[['Segment', 'Sector', 'Sector Weight (Ibovespa)', 'Sector Weight (Portfolio)',
             'Company', 'Ticker', 'Weight', 'Rating', 'Target Price']]
    df = df.sort_values(['Segment', 'Sector'], ignore_index=True)

    if idioma == 'EN':
        df['Sector'] = df['Sector'].replace(ABREV_EN)   # só abrevia
    else:
        df['Sector']  = df['Sector'].replace(SETOR_PT).replace(ABREV_PT)  # traduz e abrevia
        df['Segment'] = df['Segment'].replace(SEGMENTO_PT)                # traduz segmento
        df['Rating']  = df['Rating'].replace({'Buy': 'Compra', 'Neutral': 'Neutro', 'Sell': 'Venda'})
        df = df.rename(columns={
            'Segment': 'Segmento',
            'Sector': 'Setor',
            'Sector Weight (Ibovespa)': 'Peso do setor (Ibovespa)',
            'Sector Weight (Portfolio)': 'Peso do setor (Carteira)',
            'Company': 'Companhia',
            'Weight': 'Peso',
            'Target Price': 'Preço-Alvo',
        })
    return df

# ------------------------------------------------------------
# 9.4  Formatação da worksheet (VISUAL, com merge)
# ------------------------------------------------------------
def _merge_groups(ws, col, first, last, body_font):
    start = first
    cur = ws.cell(row=start, column=col).value
    for row in range(first + 1, last + 2):
        val = ws.cell(row=row, column=col).value if row <= last else None
        if val != cur:
            if row - 1 > start:
                ws.merge_cells(start_row=start, start_column=col, end_row=row - 1, end_column=col)
                c = ws.cell(row=start, column=col)
                c.alignment = Alignment(horizontal='center', vertical='center')
                c.font = body_font
            start = row
            if row <= last:
                cur = val


def _preencher_sheet(ws, df, titulo, idioma='EN'):
    header_font = Font(name='Roboto Light', size=7, bold=True, color='000000')
    body_font   = Font(name='Roboto Light', size=7.5, bold=False, color='000000')
    title_font  = Font(name='Roboto Light', size=7.5, bold=True, color='000000')

    if idioma == 'EN':
        header_color, border_color = '8EB3DF', '8EB3DF'
    else:
        header_color, border_color = 'FFBC00', 'D9D9D9'

    ws['A1'] = titulo
    ws['A1'].font = title_font
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=len(df.columns))

    start_row = 3
    header_fill = PatternFill('solid', fgColor=header_color)
    for ci, col in enumerate(df.columns, start=1):
        c = ws.cell(row=start_row, column=ci, value=col)
        c.font = header_font
        c.alignment = Alignment(horizontal='center', vertical='center')
        c.fill = header_fill

    for ri, row in enumerate(df.itertuples(index=False), start=start_row + 1):
        for cj, val in enumerate(row, start=1):
            c = ws.cell(row=ri, column=cj, value=val)
            c.font = body_font
            c.alignment = Alignment(horizontal='center', vertical='center')

    last_row = start_row + len(df)
    for r in range(start_row + 1, last_row + 1):
        ws.cell(row=r, column=3).number_format = '0.0%'
        ws.cell(row=r, column=4).number_format = '0.0%'
        ws.cell(row=r, column=7).number_format = '0.0%'
        ws.cell(row=r, column=9).number_format = '"R$" #,##0.00'

    thin = Side(border_style='thin', color=border_color)
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    for r in range(start_row, last_row + 1):
        for c in range(1, len(df.columns) + 1):
            ws.cell(row=r, column=c).border = border

    for col in [1, 2, 3, 4]:
        _merge_groups(ws, col, start_row + 1, last_row, body_font)

    widths = {1: 15, 2: 18, 3: 22, 4: 24, 5: 25, 6: 10, 7: 10, 8: 10, 9: 15}
    for ci, w in widths.items():
        ws.column_dimensions[chr(ord('A') + ci - 1)].width = w


def _preencher_sheet_flat(ws, df, idioma='EN'):
    """
    Versão SEM merge: repete Segment/Sector/Pesos em todas as linhas.
    Ideal para query (Power Query, Ctrl+T, fórmulas). Header na linha 1.
    """
    header_font = Font(name='Roboto Light', size=8, bold=True, color='000000')
    body_font   = Font(name='Roboto Light', size=8, bold=False, color='000000')

    header_color = '8EB3DF' if idioma == 'EN' else 'FFBC00'
    header_fill = PatternFill('solid', fgColor=header_color)

    for ci, col in enumerate(df.columns, start=1):
        c = ws.cell(row=1, column=ci, value=col)
        c.font = header_font
        c.alignment = Alignment(horizontal='center', vertical='center')
        c.fill = header_fill

    for ri, row in enumerate(df.itertuples(index=False), start=2):
        for cj, val in enumerate(row, start=1):
            c = ws.cell(row=ri, column=cj, value=val)
            c.font = body_font
            c.alignment = Alignment(horizontal='center', vertical='center')

    last_row = 1 + len(df)
    for r in range(2, last_row + 1):
        ws.cell(row=r, column=3).number_format = '0.0%'
        ws.cell(row=r, column=4).number_format = '0.0%'
        ws.cell(row=r, column=7).number_format = '0.0%'
        ws.cell(row=r, column=9).number_format = '"R$" #,##0.00'

    widths = {1: 15, 2: 18, 3: 22, 4: 24, 5: 25, 6: 10, 7: 10, 8: 10, 9: 15}
    for ci, w in widths.items():
        ws.column_dimensions[chr(ord('A') + ci - 1)].width = w

# ------------------------------------------------------------
# 9.5  Gera 1 arquivo por carteira (abas: ENG, PT, ENG_data, PT_data)
# ------------------------------------------------------------
nomes_exib = {
    'Carteira - TOP Ações XP':      'Top Ideas',
    'Carteira - TOP DIVIDENDOS XP': 'Top Dividends',
    'Carteira - TOP SMALL CAPS XP': 'Top Small Caps',
    'Carteira - ESG XP':            'Top ESG',
}
arq_comp = {
    'Carteira - TOP Ações XP':      'composicao_top_acoes',
    'Carteira - TOP DIVIDENDOS XP': 'composicao_top_dividendos',
    'Carteira - TOP SMALL CAPS XP': 'composicao_top_small_caps',
    'Carteira - ESG XP':            'composicao_esg',
}

dfs_composicao = {}
for portfolio in portfolio_names:
    nome = nomes_exib[portfolio]
    df_en = montar_df_composicao(composition_dict[portfolio], idioma='EN')
    df_pt = montar_df_composicao(composition_dict[portfolio], idioma='PT')
    dfs_composicao[portfolio] = {'EN': df_en, 'PT': df_pt}

    wb = Workbook()
    wb.remove(wb.active)

    # abas VISUAIS (com merge) — para relatório
    ws_en = wb.create_sheet(title='ENG')
    _preencher_sheet(ws_en, df_en, f"{nome} Portfolio", idioma='EN')

    ws_pt = wb.create_sheet(title='PT')
    _preencher_sheet(ws_pt, df_pt, f"Carteira {nome}", idioma='PT')

    # abas FLAT (sem merge) — para query / Power Query
    ws_en_q = wb.create_sheet(title='ENG_data')
    _preencher_sheet_flat(ws_en_q, df_en, idioma='EN')

    ws_pt_q = wb.create_sheet(title='PT_data')
    _preencher_sheet_flat(ws_pt_q, df_pt, idioma='PT')

    wb.save(f"{output_dir}\\{arq_comp[portfolio]}.xlsx")
    print(f"Composição gerada (visual + data): {nome}")


# ============================================================
# 10. GIRO MÉDIO MENSAL (turnover two-way, alvo-a-alvo)
# ============================================================
def giro_medio_mensal(composition):
    """
    Giro médio mensal (two-way): para cada rebalanceamento,
        giro = Σ_i |peso_alvo_novo_i − peso_alvo_antigo_i|
    comparando os pesos-alvo (normalizados) de rebalanceamentos consecutivos.
    Ativos que entram/saem contam com peso 0 do outro lado.
    Retorna a média dos giros ao longo dos rebalanceamentos (mensais).
    """
    comp_long = composition.melt(id_vars='cod_ativo', var_name='data_rebal', value_name='peso')
    comp_long['data_rebal'] = pd.to_datetime(comp_long['data_rebal'])
    comp_long = comp_long.dropna(subset=['peso'])

    rebal_dates = sorted(comp_long['data_rebal'].unique())
    giros = []
    for i in range(1, len(rebal_dates)):
        ant = (comp_long[comp_long['data_rebal'] == rebal_dates[i - 1]]
               .set_index('cod_ativo')['peso'])
        nov = (comp_long[comp_long['data_rebal'] == rebal_dates[i]]
               .set_index('cod_ativo')['peso'])
        ant = ant / ant.sum()
        nov = nov / nov.sum()
        ativos = ant.index.union(nov.index)
        ant = ant.reindex(ativos, fill_value=0.0)
        nov = nov.reindex(ativos, fill_value=0.0)
        giros.append(float((nov - ant).abs().sum()))   # two-way (sem 0.5)

    return float(np.mean(giros)) if giros else np.nan


# ============================================================
# 11. ESTATÍSTICAS MENSAIS (a partir da curva base 100 da carteira)
# ============================================================
def _serie_mensal(serie):
    """
    Nível base 100 no fim de cada mês (último valor observado no mês).
    Usa agrupamento por período mensal (robusto a versões do pandas — evita
    o alias de frequência 'M'/'ME' do resample, que mudou entre versões).
    """
    s = serie.dropna()
    if s.empty:
        return s
    grp = s.groupby(s.index.to_period('M')).last()
    grp.index = grp.index.to_timestamp()   # índice datetime (início do mês), monotônico
    return grp


def retornos_mensais(serie):
    """Retornos mensais (fim de mês -> fim de mês). Exclui o mês parcial de inception."""
    return _serie_mensal(serie).pct_change().dropna()


def estatisticas_carteira(serie):
    """
    Estatísticas da lâmina, sobre os retornos mensais e a curva mensal (base 100):
      Meses positivos/negativos, retorno médio/máximo/mínimo mensal,
      maior drawdown e sua duração (em meses, do topo ao fundo).
    """
    r = retornos_mensais(serie)
    m = _serie_mensal(serie)   # níveis base 100 (para o drawdown)

    if r.empty or m.empty:
        return {
            'Meses positivos': np.nan, 'Meses negativos': np.nan,
            'Retorno médio mensal': np.nan, 'Retorno máximo mensal': np.nan,
            'Retorno mínimo mensal': np.nan, 'Maior drawdown': np.nan,
            'Duração do maior drawdown (meses)': np.nan,
        }

    run_max = m.cummax()
    dd = m / run_max - 1.0
    maior_dd = float(dd.min())
    fundo = dd.idxmin()
    topo = m.loc[:fundo].idxmax()
    duracao = (fundo.year - topo.year) * 12 + (fundo.month - topo.month)

    return {
        'Meses positivos':        int((r > 0).sum()),
        'Meses negativos':        int((r < 0).sum()),
        'Retorno médio mensal':   float(r.mean()),
        'Retorno máximo mensal':  float(r.max()),
        'Retorno mínimo mensal':  float(r.min()),
        'Maior drawdown':         maior_dd,
        'Duração do maior drawdown (meses)': int(duracao),
    }


# ============================================================
# 12. TABELAS DA LÂMINA (reaproveitam _ret_mes / _ret_ano / _ret_periodo)
# ============================================================
def _nome_col_carteira(df_port):
    """Nome da coluna da carteira (a que não é benchmark) dentro do df de resultados."""
    for c in df_port.columns:
        if c not in ('Ibovespa', 'SMLL', 'ISEE'):
            return c
    return df_port.columns[0]


def tabela_performance_mensal_ano(df_port, ano=None):
    """
    Tabela 'Retornos <ano>' da lâmina: Jan..Dez do ano-calendário + total do ano (YTD).
    Linhas = carteira + benchmarks; colunas = meses + ano.
    Meses futuros ficam NaN (exibidos como '-' na lâmina).
    """
    data_fim = df_port.dropna(how='all').index.max()
    if ano is None:
        ano = data_fim.year

    ini_ano = pd.Timestamp(ano, 1, 1) - pd.Timedelta(days=1)
    meses_lbl = [pd.Timestamp(ano, mes, 1).strftime('%b') for mes in range(1, 13)]

    linhas = {}
    for col in df_port.columns:
        s = df_port[col]
        row = {}
        for mes, lbl in enumerate(meses_lbl, start=1):
            # mês ainda não iniciado (futuro) -> vazio (exibido como '-' na lâmina)
            if pd.Timestamp(ano, mes, 1) > data_fim:
                row[lbl] = np.nan
            else:
                row[lbl] = _ret_mes(s, ano, mes)
        row[str(ano)] = _ret_periodo(s, ini_ano, data_fim)
        linhas[col] = row

    tabela = pd.DataFrame(linhas).T                 # linhas = carteira/benchs
    return tabela[meses_lbl + [str(ano)]]


def tabela_retorno_anual(df_port):
    """
    'Retornos anos anteriores': retorno de cada ano-calendário fechado
    (do inception até o ano anterior ao atual). Linhas = anos; colunas = carteira + benchs.
    """
    data_fim = df_port.dropna(how='all').index.max()
    inception = df_port.index.min()
    anos = list(range(data_fim.year - 1, inception.year - 1, -1))

    linhas = {}
    for col in df_port.columns:
        s = df_port[col]
        linhas[col] = {str(a): _ret_ano(s, a) for a in anos}

    return pd.DataFrame(linhas)                     # index = anos, colunas = carteira/benchs


def tabela_retornos_acumulados(df_port):
    """
    'Retornos acumulados': 12m, 24m, 36m e desde o início (Retorno Acumulado).
    Quando a janela ultrapassa o inception, usa o inception como ponto inicial.
    Linhas = período; colunas = carteira + benchmarks.
    """
    data_fim = df_port.dropna(how='all').index.max()
    inception = df_port.index.min()

    periodos = {
        'Últimos 12 meses': data_fim - pd.DateOffset(months=12),
        'Últimos 24 meses': data_fim - pd.DateOffset(months=24),
        'Últimos 36 meses': data_fim - pd.DateOffset(months=36),
        'Retorno Acumulado': inception,
    }

    linhas = {}
    for col in df_port.columns:
        s = df_port[col]
        linhas[col] = {k: _ret_periodo(s, ini, data_fim) for k, ini in periodos.items()}

    return pd.DataFrame(linhas)                     # index = períodos, colunas = carteira/benchs


def info_adicionais_lamina(df_port, composition, serie_cdi):
    """
    Bloco 'Informações adicionais' pedido: giro médio mensal, volatilidade
    anualizada (12m) e Índice de Sharpe (12m). Vol/Sharpe vêm de indicadores_12m.
    """
    ind = indicadores_12m(df_port, serie_cdi)
    nome_cart = _nome_col_carteira(df_port)
    return {
        'Giro médio mensal':        giro_medio_mensal(composition),
        'Volatilidade anualizada':  float(ind.loc['Volatilidade', nome_cart]),
        'Índice de Sharpe':         float(ind.loc['Sharpe', nome_cart]),
    }


# ============================================================
# 13. EXPORTA 1 EXCEL POR CARTEIRA
#     Cada aba replica EXATAMENTE a tabela correspondente do PPT
#     (mesmos rótulos, ordem e formatação pt-BR), para copiar e colar
#     os valores direto na lâmina. Valores gravados como TEXTO já formatado.
# ============================================================
arq_lamina = {
    'Carteira - TOP Ações XP':      'lamina_dados_top_acoes',
    'Carteira - TOP DIVIDENDOS XP': 'lamina_dados_top_dividendos',
    'Carteira - TOP SMALL CAPS XP': 'lamina_dados_top_small_caps',
    'Carteira - ESG XP':            'lamina_dados_esg',
}

# meses em português (para o cabeçalho da tabela "Retornos <ano>")
MESES_PT = ['Jan', 'Fev', 'Mar', 'Abr', 'Mai', 'Jun',
            'Jul', 'Ago', 'Set', 'Out', 'Nov', 'Dez']

# rótulo da linha da carteira na tabela mensal (igual ao PPT)
rotulo_carteira_pt = {
    'Carteira - TOP Ações XP':      'Top Ações',
    'Carteira - TOP DIVIDENDOS XP': 'Top Dividendos',
    'Carteira - TOP SMALL CAPS XP': 'Top Small Caps',
    'Carteira - ESG XP':            'ESG',
}

# ------------------------------------------------------------------
# Valores FIXOS de "Informações adicionais" (não são calculados aqui).
# >>> ATUALIZAR MENSALMENTE <<< — deixe '' para preencher à mão na lâmina.
# Pré-preenchido com os valores da lâmina de Junho/2026 (Top Ações).
# ------------------------------------------------------------------
info_fixas = {
    'Carteira - TOP Ações XP': {
        'corretagem':     '1,41%',
        'qtd_clientes':   '35 mil',
        'media_clientes': '32 mil',
    },
    'Carteira - TOP DIVIDENDOS XP': {'corretagem': '', 'qtd_clientes': '', 'media_clientes': ''},
    'Carteira - TOP SMALL CAPS XP': {'corretagem': '', 'qtd_clientes': '', 'media_clientes': ''},
    'Carteira - ESG XP':            {'corretagem': '', 'qtd_clientes': '', 'media_clientes': ''},
}


# ------------------------------------------------------------------
# Formatação pt-BR (texto) — igual ao que aparece no PPT
# ------------------------------------------------------------------
def _eh_nan(x):
    return x is None or (isinstance(x, float) and np.isnan(x))


def _fmt_pct_br(x, dec=1, dash='-'):
    """0.123 -> '12,3%' ; negativo -> '-6,7%' ; NaN -> '-'."""
    if _eh_nan(x):
        return dash
    return f"{x * 100:.{dec}f}%".replace('.', ',')


def _fmt_num_br(x, dec=2, dash='-'):
    """0.46 -> '0,46' ; NaN -> '-'."""
    if _eh_nan(x):
        return dash
    return f"{x:.{dec}f}".replace('.', ',')


def _fmt_int_br(x, dash='-'):
    if _eh_nan(x):
        return dash
    return str(int(round(x)))


def _fmt_dur(n, dash='-'):
    if _eh_nan(n):
        return dash
    n = int(round(n))
    return f"{n} mês" if n == 1 else f"{n} meses"


def _bench_lbl(cod):
    """Rótulo do benchmark nas tabelas anuais/acumuladas (PPT usa 'IBOV')."""
    return 'IBOV' if cod == 'Ibovespa' else cod


# ------------------------------------------------------------------
# Construtores das grades (list-of-lists) idênticas ao PPT
# ------------------------------------------------------------------
def _grid_performance_mensal(df_port, portfolio, ano=None):
    """Tabela 'Retornos <ano>': linha da carteira + benchmarks, meses Jan..Dez + ano."""
    t = tabela_performance_mensal_ano(df_port, ano)      # linhas = cart/bench ; cols = [Jan..Dec(en), ano]
    ano_lbl = t.columns[-1]
    meses_en = list(t.columns[:-1])

    nome_cart = _nome_col_carteira(df_port)
    ordem = [nome_cart] + [c for c in df_port.columns if c != nome_cart]

    grid = [[''] + MESES_PT + [str(ano_lbl)]]
    for col in ordem:
        rot = rotulo_carteira_pt.get(portfolio, nome_cart) if col == nome_cart else col
        vals = t.loc[col]
        linha = [rot] + [_fmt_pct_br(vals[m], 1) for m in meses_en] + [_fmt_pct_br(vals[ano_lbl], 1)]
        grid.append(linha)
    return grid


def _grid_retorno_anual(df_port):
    """Tabela 'Retornos anos anteriores': linhas = anos ; colunas = Carteira / IBOV."""
    t = tabela_retorno_anual(df_port)                    # linhas = anos ; cols = cart/bench
    nome_cart = _nome_col_carteira(df_port)
    benches = [c for c in df_port.columns if c != nome_cart]

    grid = [['Retornos anos anteriores', 'Carteira'] + [_bench_lbl(b) for b in benches]]
    for ano in t.index:
        linha = [str(ano), _fmt_pct_br(t.loc[ano, nome_cart], 1)]
        linha += [_fmt_pct_br(t.loc[ano, b], 1) for b in benches]
        grid.append(linha)
    return grid


def _grid_retornos_acumulados(df_port):
    """Tabela 'Retornos acumulados': linhas = períodos ; colunas = Carteira / IBOV."""
    t = tabela_retornos_acumulados(df_port)              # linhas = períodos ; cols = cart/bench
    nome_cart = _nome_col_carteira(df_port)
    benches = [c for c in df_port.columns if c != nome_cart]

    grid = [['Retornos acumulados', 'Carteira'] + [_bench_lbl(b) for b in benches]]
    for per in t.index:
        linha = [per, _fmt_pct_br(t.loc[per, nome_cart], 1)]
        linha += [_fmt_pct_br(t.loc[per, b], 1) for b in benches]
        grid.append(linha)
    return grid


def _grid_estatisticas(serie):
    """Tabela 'Estatísticas' (2 colunas), igual ao PPT."""
    e = estatisticas_carteira(serie)
    return [
        ['Estatísticas', ''],
        ['Meses positivos',        _fmt_int_br(e['Meses positivos'])],
        ['Meses negativos',        _fmt_int_br(e['Meses negativos'])],
        ['Retorno médio mensal',   _fmt_pct_br(e['Retorno médio mensal'], 2)],
        ['Retorno máximo mensal',  _fmt_pct_br(e['Retorno máximo mensal'], 2)],
        ['Retorno mínimo mensal',  _fmt_pct_br(e['Retorno mínimo mensal'], 2)],
        ['Maior drawdown',         _fmt_pct_br(e['Maior drawdown'], 2)],
        ['Duração do maior drawdown', _fmt_dur(e['Duração do maior drawdown (meses)'])],
    ]


def _grid_info_adicionais(df_port, composition, serie_cdi, portfolio):
    """Tabela 'Informações adicionais' (2 colunas). Giro/Vol/Sharpe calculados; resto fixo."""
    info = info_adicionais_lamina(df_port, composition, serie_cdi)
    fix = info_fixas.get(portfolio, {})
    return [
        ['Informações adicionais', ''],
        ['Giro médio mensal*',            _fmt_pct_br(info['Giro médio mensal'], 1)],
        ['Corretagem estimada anual**',   fix.get('corretagem', '')],
        ['Volatilidade anualizada*',      _fmt_pct_br(info['Volatilidade anualizada'], 2)],
        ['Índice de Sharpe*',             _fmt_num_br(info['Índice de Sharpe'], 2)],
        ['Quantidade de clientes',        fix.get('qtd_clientes', '')],
        ['Média de clientes na Carteira por mês*', fix.get('media_clientes', '')],
    ]


def _escreve_grid(ws, formats, grid, larg_primeira=30, larg_demais=11):
    """Grava uma grade (list-of-lists de strings) na worksheet, com a 1ª linha
    e a 1ª coluna em destaque (igual ao layout de tabela do PPT)."""
    ncols = max(len(r) for r in grid)
    for r, row in enumerate(grid):
        for c, val in enumerate(row):
            if r == 0:
                ws.write(r, c, val, formats['hdr'])
            elif c == 0:
                ws.write(r, c, val, formats['lbl'])
            else:
                ws.write(r, c, val, formats['cell'])
    ws.set_column(0, 0, larg_primeira)
    if ncols > 1:
        ws.set_column(1, ncols - 1, larg_demais)


def exportar_lamina(df_port, composition, serie_cdi, caminho, portfolio):
    """Gera o Excel da lâmina: 1 aba por gráfico/tabela, replicando o PPT."""
    nome_cart = _nome_col_carteira(df_port)

    with pd.ExcelWriter(caminho, engine='xlsxwriter') as writer:
        wb = writer.book
        formats = {
            'data':     wb.add_format({'num_format': 'dd/mm/yyyy', 'align': 'center'}),
            'num2b100': wb.add_format({'num_format': '0.00', 'align': 'center'}),
            'hdr':      wb.add_format({'bold': True, 'bg_color': '#1F3864', 'font_color': 'white',
                                       'align': 'center', 'border': 1}),
            'lbl':      wb.add_format({'bold': True, 'align': 'left', 'border': 1}),
            'cell':     wb.add_format({'align': 'center', 'border': 1}),
        }

        # aba 1 — base100 (dados do gráfico "Performance desde o início vs. Ibovespa")
        df_chart = df_port.rename(columns={nome_cart: rotulo_carteira_pt.get(portfolio, nome_cart)})
        df_chart.to_excel(writer, sheet_name='base100')
        ws = writer.sheets['base100']
        ws.set_column(0, 0, 12, formats['data'])
        ws.set_column(1, len(df_chart.columns), 16, formats['num2b100'])

        # aba 2 — Retornos <ano> (performance mensal)
        ws = wb.add_worksheet('performance_mensal')
        _escreve_grid(ws, formats, _grid_performance_mensal(df_port, portfolio),
                      larg_primeira=14, larg_demais=8)

        # aba 3 — Retornos anos anteriores
        ws = wb.add_worksheet('retorno_anual')
        _escreve_grid(ws, formats, _grid_retorno_anual(df_port),
                      larg_primeira=24, larg_demais=11)

        # aba 4 — Retornos acumulados
        ws = wb.add_worksheet('retornos_acumulados')
        _escreve_grid(ws, formats, _grid_retornos_acumulados(df_port),
                      larg_primeira=20, larg_demais=11)

        # aba 5 — Estatísticas
        ws = wb.add_worksheet('estatisticas')
        _escreve_grid(ws, formats, _grid_estatisticas(df_port[nome_cart]),
                      larg_primeira=28, larg_demais=12)

        # aba 6 — Informações adicionais
        ws = wb.add_worksheet('info_adicionais')
        _escreve_grid(ws, formats, _grid_info_adicionais(df_port, composition, serie_cdi, portfolio),
                      larg_primeira=36, larg_demais=12)

    print(f"Lâmina (dados) exportada: {caminho}")


for portfolio in portfolio_names:
    exportar_lamina(
        df_port=_df_para_lamina(portfolio),          
        composition=composition_dict[portfolio],
        serie_cdi=cdi,
        caminho=f"{output_dir}\\{arq_lamina[portfolio]}.xlsx",
        portfolio=portfolio,
    )


# ============================================================
# 14. ATUALIZA O PPT DA LÂMINA (preenche tabelas + gráfico no lugar,
#     mantendo 100% da formatação original — fontes, cores, bordas, gráfico)
# ============================================================
import os
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

# ------------------------------------------------------------------
# Caminhos: template (.pptx base do mês) e saída de cada carteira.
# >>> Preencha 'template' e 'saida' para cada carteira que quiser gerar. <<<
# Só as entradas com 'template' existente são processadas (as demais são puladas).
# ------------------------------------------------------------------
ppt_config = {
    'Carteira - TOP Ações XP': {
        'template': r"\\xpdocs\Research\Equities\Estrategia\Carteiras\Carteiras de Ações XP\Templates\Lâmina Comercial - Top Ações - Junho 2026.pptx",
        'saida':    r"\\xpdocs\Research\Equities\Estrategia\Carteiras\Carteiras de Ações XP\Lâmina Comercial\Lâmina Comercial - Top Ações.pptx",
    },
    'Carteira - TOP DIVIDENDOS XP': {'template': r"\\xpdocs\Research\Equities\Estrategia\Carteiras\Carteiras de Ações XP\Templates\Lâmina Comercial - Top Dividendos - Junho 2026.pptx",
                                     'saida': r"\\xpdocs\Research\Equities\Estrategia\Carteiras\Carteiras de Ações XP\Lâmina Comercial\Lâmina Comercial - Top Dividendos.pptx"},

    'Carteira - TOP SMALL CAPS XP': {'template': r"\\xpdocs\Research\Equities\Estrategia\Carteiras\Carteiras de Ações XP\Templates\Lâmina Comercial - Top Small Caps - Junho 2026.pptx",
                                     'saida': r"\\xpdocs\Research\Equities\Estrategia\Carteiras\Carteiras de Ações XP\Lâmina Comercial\Lâmina Comercial - Top Small Caps.pptx"},
}


def _set_cell_text(cell, value):
    """Troca o texto da célula preservando a formatação do primeiro run."""
    p = cell.text_frame.paragraphs[0]
    runs = p.runs
    if runs:
        runs[0].text = str(value)
        for r in runs[1:]:                 # remove runs extras (mantém a fonte do 1º)
            r._r.getparent().remove(r._r)
    else:
        p.add_run().text = str(value)


def _classifica_tabelas(slide):
    """Localiza cada tabela da lâmina pelo texto do cabeçalho (robusto a posição)."""
    achadas = {}
    for shp in slide.shapes:
        if not shp.has_table:
            continue
        t = shp.table
        r0c0 = t.rows[0].cells[0].text.strip()
        if r0c0 == 'Retornos anos anteriores':
            achadas['anos'] = t
        elif r0c0 == 'Retornos acumulados':
            achadas['acumulados'] = t
        elif r0c0 == 'Estatísticas':
            achadas['estatisticas'] = t
        elif r0c0 == 'Informações adicionais':
            achadas['info'] = t
        else:
            # tabela mensal "Retornos <ano>": cabeçalho começa com 'Jan'
            try:
                if t.rows[0].cells[1].text.strip() == 'Jan':
                    achadas['mensal'] = t
            except (IndexError, AttributeError):
                pass
    return achadas


def _mapa_colunas(table, nome_cart):
    """Mapeia índice de coluna -> série, lendo o cabeçalho (Carteira / IBOV / SMLL / ISEE)."""
    header = [c.text.strip() for c in table.rows[0].cells]
    col_map = {}
    for ci, h in enumerate(header):
        if ci == 0:
            continue
        if h == 'Carteira':
            col_map[ci] = nome_cart
        elif h == 'IBOV':
            col_map[ci] = 'Ibovespa'
        elif h in ('SMLL', 'ISEE'):
            col_map[ci] = h
    return col_map


def _atualiza_tabela_mensal(table, df_port, portfolio, ano=None):
    t = tabela_performance_mensal_ano(df_port, ano)      # linhas cart/bench ; cols Jan..Dez(en)+ano
    ano_lbl = t.columns[-1]
    meses_en = list(t.columns[:-1])
    nome_cart = _nome_col_carteira(df_port)
    rot_cart = rotulo_carteira_pt.get(portfolio, nome_cart)

    # atualiza o rótulo do ano no cabeçalho (última coluna da 1ª linha)
    cabec = table.rows[0].cells
    _set_cell_text(cabec[len(cabec) - 1], str(ano_lbl))

    for row in list(table.rows)[1:]:
        rot = row.cells[0].text.strip()
        if rot in (rot_cart, nome_cart):
            serie_col = nome_cart
        elif rot in df_port.columns:
            serie_col = rot
        elif rot == 'Ibovespa' and 'Ibovespa' in df_port.columns:
            serie_col = 'Ibovespa'
        else:
            continue
        vals = t.loc[serie_col]
        cells = row.cells
        for j, m in enumerate(meses_en, start=1):        # c1..c12 = meses
            _set_cell_text(cells[j], _fmt_pct_br(vals[m], 1))
        _set_cell_text(cells[len(cells) - 1], _fmt_pct_br(vals[ano_lbl], 1))   # última = ano


def _atualiza_tabela_anos(table, df_port):
    nome_cart = _nome_col_carteira(df_port)
    col_map = _mapa_colunas(table, nome_cart)
    for row in list(table.rows)[1:]:
        try:
            ano = int(row.cells[0].text.strip())
        except ValueError:
            continue
        for ci, serie in col_map.items():
            if serie in df_port.columns:
                _set_cell_text(row.cells[ci], _fmt_pct_br(_ret_ano(df_port[serie], ano), 1))


def _atualiza_tabela_acumulados(table, df_port):
    nome_cart = _nome_col_carteira(df_port)
    t = tabela_retornos_acumulados(df_port)              # linhas períodos ; cols cart/bench
    col_map = _mapa_colunas(table, nome_cart)
    for row in list(table.rows)[1:]:
        per = row.cells[0].text.strip()
        if per not in t.index:
            continue
        for ci, serie in col_map.items():
            if serie in t.columns:
                _set_cell_text(row.cells[ci], _fmt_pct_br(t.loc[per, serie], 1))


def _atualiza_2col_por_rotulo(table, mapa):
    """Atualiza a 2ª coluna das linhas cujo rótulo (1ª coluna) casa com o mapa."""
    for row in table.rows:
        rot = row.cells[0].text.strip()
        if rot in mapa and mapa[rot] is not None:
            _set_cell_text(row.cells[1], mapa[rot])


def _mapa_estatisticas(serie):
    e = estatisticas_carteira(serie)
    return {
        'Meses positivos':        _fmt_int_br(e['Meses positivos']),
        'Meses negativos':        _fmt_int_br(e['Meses negativos']),
        'Retorno médio mensal':   _fmt_pct_br(e['Retorno médio mensal'], 2),
        'Retorno máximo mensal':  _fmt_pct_br(e['Retorno máximo mensal'], 2),
        'Retorno mínimo mensal':  _fmt_pct_br(e['Retorno mínimo mensal'], 2),
        'Maior drawdown':         _fmt_pct_br(e['Maior drawdown'], 2),
        'Duração do maior drawdown': _fmt_dur(e['Duração do maior drawdown (meses)']),
    }


def _mapa_info(df_port, composition, serie_cdi):
    """Só os 3 campos calculados. Os fixos (corretagem, clientes) ficam intocados no PPT."""
    info = info_adicionais_lamina(df_port, composition, serie_cdi)
    return {
        'Giro médio mensal*':       _fmt_pct_br(info['Giro médio mensal'], 1),
        'Volatilidade anualizada*': _fmt_pct_br(info['Volatilidade anualizada'], 2),
        'Índice de Sharpe*':        _fmt_num_br(info['Índice de Sharpe'], 2),
    }


def _nome_serie_xml(serie_element):
    """Lê o nome original da série no XML (c:tx), para preservar a legenda."""
    tx = serie_element.find(qn('c:tx'))
    if tx is None:
        return None
    valores = [v.text for v in tx.iter(qn('c:v')) if v.text]
    return valores[0] if valores else None


def _atualiza_grafico(shape, df_port):
    """Substitui os dados do gráfico de linha, preservando estilo E nomes das séries."""
    nome_cart = _nome_col_carteira(df_port)
    chart = shape.chart

    # ordem das colunas = mesma ordem das séries do gráfico (carteira, depois benchmarks)
    cols = [nome_cart] + [c for c in df_port.columns if c != nome_cart]
    nomes_orig = [_nome_serie_xml(s._element) for s in chart.series]

    def _limpa(col):
        return [None if pd.isna(v) else float(v) for v in df_port[col].tolist()]

    cd = CategoryChartData(number_format=r'[$-416]mmm\-yy;@')
    cd.categories = [pd.Timestamp(d).to_pydatetime() for d in df_port.index]
    for i, col in enumerate(cols):
        nome = nomes_orig[i] if i < len(nomes_orig) and nomes_orig[i] else col
        cd.add_series(nome, _limpa(col))
    chart.replace_data(cd)


def atualizar_ppt(caminho_template, caminho_saida, df_port, composition, serie_cdi, portfolio, ano=None):
    """Abre a lâmina, preenche tabelas e gráfico, e salva mantendo a formatação."""
    prs = Presentation(caminho_template)
    for slide in prs.slides:
        tabs = _classifica_tabelas(slide)
        if 'mensal' in tabs:
            _atualiza_tabela_mensal(tabs['mensal'], df_port, portfolio, ano)
        if 'anos' in tabs:
            _atualiza_tabela_anos(tabs['anos'], df_port)
        if 'acumulados' in tabs:
            _atualiza_tabela_acumulados(tabs['acumulados'], df_port)
        if 'estatisticas' in tabs:
            _atualiza_2col_por_rotulo(
                tabs['estatisticas'], _mapa_estatisticas(df_port[_nome_col_carteira(df_port)]))
        if 'info' in tabs:
            _atualiza_2col_por_rotulo(tabs['info'], _mapa_info(df_port, composition, serie_cdi))
        for shp in slide.shapes:
            if shp.has_chart:
                _atualiza_grafico(shp, df_port)
    prs.save(caminho_saida)
    print(f"PPT atualizado: {caminho_saida}")


for portfolio, cfg in ppt_config.items():
    template = cfg.get('template', '')
    if not template or not os.path.exists(template):
        print(f"[PULADO] template não encontrado para {portfolio}: {template}")
        continue
    atualizar_ppt(
        caminho_template=template,
        caminho_saida=cfg['saida'],
        df_port=_df_para_lamina(portfolio),         
        composition=composition_dict[portfolio],
        serie_cdi=cdi,
        portfolio=portfolio,
    )


# ============================================================
# 15. DECOMPOSIÇÃO DE RETORNO (RETURN ATTRIBUTION) POR PAPEL
#     contribuição_i = peso_i * retorno_i  (mês atual e mês anterior)
#     Soma das contribuições = retorno da carteira no mês.
# ============================================================
def _retorno_papel_mes(cod_ativo, ano, mes):
    """Retorno do papel no mês cheio (fim do mês anterior -> fim do mês),
    usando o último preço ajustado disponível <= cada data de corte."""
    ini = pd.Timestamp(ano, mes, 1) - pd.Timedelta(days=1)
    fim = pd.Timestamp(ano, mes, 1) + pd.offsets.MonthEnd(0)
    p_ini = _preco_em(cod_ativo, ini)
    p_fim = _preco_em(cod_ativo, fim)
    if not p_ini or not p_fim or np.isnan(p_ini) or np.isnan(p_fim):
        return np.nan
    return p_fim / p_ini - 1


def _rebal_vigente_no_mes(comp_long, rebal_dates, ano, mes):
    """Rebalanceamento vigente durante o mês (último rebal <= fim do mês)."""
    fim_mes = pd.Timestamp(ano, mes, 1) + pd.offsets.MonthEnd(0)
    vigentes = [d for d in rebal_dates if pd.Timestamp(d) <= fim_mes]
    return vigentes[-1] if vigentes else rebal_dates[0]


def decomposicao_retorno(composition, ano, mes):
    """
    Return attribution do mês (ano, mes):
      - peso: peso-alvo normalizado do rebal vigente no mês
      - retorno: retorno mensal do papel
      - contribuição: peso * retorno
    A linha 'Carteira (total)' traz a soma das contribuições (retorno do mês).
    """
    comp_long = composition.melt(id_vars='cod_ativo', var_name='data_rebal', value_name='peso')
    comp_long['data_rebal'] = pd.to_datetime(comp_long['data_rebal'])
    comp_long = comp_long.dropna(subset=['peso'])
    rebal_dates = sorted(comp_long['data_rebal'].unique())

    rebal_ref = _rebal_vigente_no_mes(comp_long, rebal_dates, ano, mes)
    comp_ref = comp_long[comp_long['data_rebal'] == rebal_ref][['cod_ativo', 'peso']].copy()
    soma_peso = comp_ref['peso'].sum()

    linhas = []
    for _, r in comp_ref.iterrows():
        cod = r['cod_ativo']
        peso = r['peso'] / soma_peso if soma_peso else np.nan
        ret = _retorno_papel_mes(cod, ano, mes)
        contrib = peso * ret if (peso is not None and not np.isnan(peso) and not np.isnan(ret)) else np.nan
        linhas.append({
            'Companhia':        mapa_nome.get(cod, ''),
            'Ticker':           cod,
            'Setor':            mapa_setor.get(cod, ''),
            'Peso':             peso,
            'Retorno no mês':   ret,
            'Contribuição':     contrib,
        })

    tab = pd.DataFrame(linhas).sort_values('Contribuição', ascending=False).reset_index(drop=True)

    # linha de total (retorno da carteira no mês = soma das contribuições)
    total = pd.DataFrame([{
        'Companhia': 'Carteira (total)',
        'Ticker': '',
        'Setor': '',
        'Peso': tab['Peso'].sum(skipna=True),
        'Retorno no mês': np.nan,
        'Contribuição': tab['Contribuição'].sum(skipna=True),
    }])
    return pd.concat([tab, total], ignore_index=True)


def gerar_decomposicoes(composition):
    """Return attribution do mês atual (aberto) e do mês anterior (fechado)."""
    data_max = md['data'].max()

    ano_atual, mes_atual = data_max.year, data_max.month

    ref_ant = pd.Timestamp(ano_atual, mes_atual, 1) - pd.Timedelta(days=1)
    ano_ant, mes_ant = ref_ant.year, ref_ant.month

    dec_atual = decomposicao_retorno(composition, ano_atual, mes_atual)
    dec_ant   = decomposicao_retorno(composition, ano_ant,   mes_ant)
    return dec_atual, dec_ant


# ------------------------------------------------------------
# 15b. EXPORTA A DECOMPOSIÇÃO (1 arquivo por carteira, 2 abas)
# ------------------------------------------------------------
arq_decomp = {
    'Carteira - TOP Ações XP':      'decomposicao_top_acoes',
    'Carteira - TOP DIVIDENDOS XP': 'decomposicao_top_dividendos',
    'Carteira - TOP SMALL CAPS XP': 'decomposicao_top_small_caps',
    'Carteira - ESG XP':            'decomposicao_esg',
}

MES_LBL_PT = {1: 'Jan', 2: 'Fev', 3: 'Mar', 4: 'Abr', 5: 'Mai', 6: 'Jun',
              7: 'Jul', 8: 'Ago', 9: 'Set', 10: 'Out', 11: 'Nov', 12: 'Dez'}

cols_pct_decomp = ['Peso', 'Retorno no mês', 'Contribuição']

for portfolio in portfolio_names:
    dec_atual, dec_ant = gerar_decomposicoes(composition_dict[portfolio])
    caminho = f"{output_dir}\\{arq_decomp[portfolio]}.xlsx"

    data_max = md['data'].max()
    lbl_atual = f"{MES_LBL_PT[data_max.month]}-{str(data_max.year)[2:]}"
    ref_ant = pd.Timestamp(data_max.year, data_max.month, 1) - pd.Timedelta(days=1)
    lbl_ant = f"{MES_LBL_PT[ref_ant.month]}-{str(ref_ant.year)[2:]}"

    with pd.ExcelWriter(caminho, engine='xlsxwriter') as writer:
        wb = writer.book
        fmt_pct = wb.add_format({'num_format': '0.0%'})

        for sheet, tab in [(f"mes_atual_{lbl_atual}", dec_atual),
                           (f"mes_anterior_{lbl_ant}", dec_ant)]:
            tab.to_excel(writer, index=False, sheet_name=sheet)
            ws = writer.sheets[sheet]
            for c in cols_pct_decomp:
                col_idx = tab.columns.get_loc(c)
                ws.set_column(col_idx, col_idx, 14, fmt_pct)
            ws.set_column(0, 0, 24)  # Companhia
            ws.set_column(2, 2, 22)  # Setor

    print(f"Decomposição de retorno exportada: {caminho}")


# ============================================================
# 16. COMPOSIÇÃO COMPACTA (Setor / Companhia / Ticker / Peso)
#     Layout igual à imagem: agrupada por setor, com merge do setor.
# ============================================================
def montar_df_composicao_compacta(composition, idioma='PT'):
    """Setor / Companhia / Ticker / Peso do último rebal, ordenado por segmento/setor."""
    comp = _composicao_vigente(composition)                        # Ticker, Weight
    comp = comp.merge(comp_sheet[['Ticker', 'Company']], on='Ticker', how='left')
    comp = comp.merge(sector_classification_raw, on='Ticker', how='left')  # Sector (GICS)
    comp = comp.merge(segment_sector, on='Sector', how='left')            # Segment (p/ ordenar)

    # ordena por segmento e setor (mesma lógica da composição completa)
    comp = comp.sort_values(['Segment', 'Sector', 'Weight'],
                            ascending=[True, True, False], ignore_index=True)

    df = comp[['Sector', 'Company', 'Ticker', 'Weight']].copy()

    if idioma == 'PT':
        df['Sector'] = df['Sector'].replace(SETOR_PT).replace(ABREV_PT)
        df = df.rename(columns={
            'Sector': 'Setor', 'Company': 'Companhia', 'Weight': 'Peso'
        })
    else:  # EN
        df['Sector'] = df['Sector'].replace(ABREV_EN)
        df = df.rename(columns={
            'Sector': 'Sector', 'Company': 'Company', 'Weight': 'Weight'
        })
    return df


def _preencher_sheet_compacta(ws, df, idioma='PT'):
    """Layout compacto da imagem: header amarelo, merge do setor por grupo."""
    header_font = Font(name='Roboto Light', size=9, bold=True, color='000000')
    body_font   = Font(name='Roboto Light', size=9, bold=False, color='000000')

    header_color = 'FFBC00' if idioma == 'PT' else '8EB3DF'
    border_color = 'D9D9D9'
    header_fill = PatternFill('solid', fgColor=header_color)

    # cabeçalho na linha 1
    for ci, col in enumerate(df.columns, start=1):
        c = ws.cell(row=1, column=ci, value=col)
        c.font = header_font
        c.alignment = Alignment(horizontal='center', vertical='center')
        c.fill = header_fill

    # corpo a partir da linha 2
    for ri, row in enumerate(df.itertuples(index=False), start=2):
        for cj, val in enumerate(row, start=1):
            c = ws.cell(row=ri, column=cj, value=val)
            c.font = body_font
            c.alignment = Alignment(horizontal='center', vertical='center')

    last_row = 1 + len(df)

    # coluna Peso (4ª) em %
    for r in range(2, last_row + 1):
        ws.cell(row=r, column=4).number_format = '0.0%'

    # bordas
    thin = Side(border_style='thin', color=border_color)
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    for r in range(1, last_row + 1):
        for c in range(1, len(df.columns) + 1):
            ws.cell(row=r, column=c).border = border

    # merge da 1ª coluna (Setor) por grupos consecutivos
    _merge_groups(ws, 1, 2, last_row, body_font)

    widths = {1: 20, 2: 26, 3: 12, 4: 10}
    for ci, w in widths.items():
        ws.column_dimensions[chr(ord('A') + ci - 1)].width = w


arq_comp_compacta = {
    'Carteira - TOP Ações XP':      'composicao_compacta_top_acoes',
    'Carteira - TOP DIVIDENDOS XP': 'composicao_compacta_top_dividendos',
    'Carteira - TOP SMALL CAPS XP': 'composicao_compacta_top_small_caps',
    'Carteira - ESG XP':            'composicao_compacta_esg',
}

for portfolio in portfolio_names:
    df_pt = montar_df_composicao_compacta(composition_dict[portfolio], idioma='PT')
    df_en = montar_df_composicao_compacta(composition_dict[portfolio], idioma='EN')

    wb = Workbook()
    wb.remove(wb.active)

    ws_pt = wb.create_sheet(title='PT')
    _preencher_sheet_compacta(ws_pt, df_pt, idioma='PT')

    ws_en = wb.create_sheet(title='ENG')
    _preencher_sheet_compacta(ws_en, df_en, idioma='EN')

    # aba flat (sem merge) para query
    ws_pt_q = wb.create_sheet(title='PT_data')
    for ci, col in enumerate(df_pt.columns, start=1):
        ws_pt_q.cell(row=1, column=ci, value=col)
    for ri, row in enumerate(df_pt.itertuples(index=False), start=2):
        for cj, val in enumerate(row, start=1):
            ws_pt_q.cell(row=ri, column=cj, value=val)
    for r in range(2, 2 + len(df_pt)):
        ws_pt_q.cell(row=r, column=4).number_format = '0.0%'

    wb.save(f"{output_dir}\\{arq_comp_compacta[portfolio]}.xlsx")
    print(f"Composição compacta gerada: {nomes_exib[portfolio]}")

# ============================================================
# 17. ATUALIZA O PPT DA "PRESTAÇÃO DE CONTAS"
#     (relatório mensal em retrato: resumo de retornos, gráfico
#      base 100, decomposição do retorno por papel e composição)
#     Reaproveita as funções deste pipeline; só adiciona os
#     preenchedores específicos deste template.
# ============================================================
from copy import deepcopy
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml import parse_xml as _parse_xml_c

MESES_EXT_PT = ['janeiro', 'fevereiro', 'março', 'abril', 'maio', 'junho',
                'julho', 'agosto', 'setembro', 'outubro', 'novembro', 'dezembro']
_C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"


def _set_para_text(p, value):
    """Troca o texto de um parágrafo preservando a formatação do 1º run."""
    runs = p.runs
    if runs:
        runs[0].text = str(value)
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p.add_run().text = str(value)


def _mes_ext(ano, mes):
    return f"{MESES_EXT_PT[mes - 1].capitalize()} de {ano}"


def _prox_mes(ano, mes):
    return (ano + (mes // 12), (mes % 12) + 1)


def _pc_ajusta_linhas(table, n):
    tbl = table._tbl
    while len(tbl.tr_lst) - 1 < n:
        tbl.append(deepcopy(tbl.tr_lst[-1]))
    while len(tbl.tr_lst) - 1 > n:
        tbl.remove(tbl.tr_lst[-1])


def _pc_unmerge_all(table):
    for r in range(len(table.rows)):
        for c in range(len(table.columns)):
            if table.cell(r, c).is_merge_origin:
                table.cell(r, c).split()


def _pc_tabela_resumo(table, df_port, ano_ref, mes_ref):
    """Tabela 3x4: Desempenho | <mês> | Acumulado <ano> | Últimos 12 meses."""
    nome_cart = _nome_col_carteira(df_port)
    data_fim = df_port.dropna(how='all').index.max()
    ini_12m = data_fim - pd.DateOffset(months=12)
    hdr = table.rows[0].cells
    _set_cell_text(hdr[1], f"{MESES_EXT_PT[mes_ref - 1].capitalize()} {ano_ref}")
    _set_cell_text(hdr[2], f"Acumulado {ano_ref}")
    for row in list(table.rows)[1:]:
        rot = row.cells[0].text.strip()
        if 'carteira' in rot.lower():
            serie = nome_cart
        else:
            serie = next((c for c in df_port.columns
                          if c == rot or _bench_lbl(c) == rot), None)
        if serie is None or serie not in df_port.columns:
            continue
        s = df_port[serie]
        _set_cell_text(row.cells[1], _fmt_pct_br(_ret_mes(s, ano_ref, mes_ref), 1))
        _set_cell_text(row.cells[2], _fmt_pct_br(_ret_ano(s, ano_ref), 1))
        _set_cell_text(row.cells[3], _fmt_pct_br(_ret_periodo(s, ini_12m, data_fim), 1))


def _pc_tabela_composicao(table, df_comp):
    """Setor | Companhia | Ticker | Peso, com o Setor mesclado por grupo."""
    _pc_unmerge_all(table)
    n = len(df_comp)
    _pc_ajusta_linhas(table, n)
    for i, r in enumerate(df_comp.itertuples(index=False), start=1):
        _set_cell_text(table.cell(i, 0), str(r.Setor))
        _set_cell_text(table.cell(i, 1), str(r.Companhia))
        _set_cell_text(table.cell(i, 2), str(r.Ticker))
        _set_cell_text(table.cell(i, 3), _fmt_pct_br(r.Peso, 1))
    r = 1
    while r <= n:
        start, val = r, table.cell(r, 0).text.strip()
        r += 1
        while r <= n and table.cell(r, 0).text.strip() == val:
            r += 1
        if r - 1 > start:
            for k in range(start + 1, r):
                _set_cell_text(table.cell(k, 0), "")
            table.cell(start, 0).merge(table.cell(r - 1, 0))


# Cores das barras do waterfall (alta / baixa / total). Amarelo XP no total.
COR_ALTA, COR_BAIXA, COR_TOTAL = '70AD47', 'C0504D', 'FFBC00'


def _waterfall_arrays(tickers, contribs, total, rotulo_total='Carteira'):
    """Vetores de um waterfall com colunas empilhadas. Uma série 'base'
    invisível flutua cada barra até o acumulado; alta/baixa desenham a
    variação e a última barra é o total. Para a base não cruzar o zero,
    ordena as contribuições de modo que o caminho fique sempre do mesmo lado
    (positivas primeiro se o total é >= 0; negativas primeiro caso contrário)."""
    total = float(total) if total is not None and not (isinstance(total, float) and np.isnan(total)) \
        else float(np.nansum([c for c in contribs if c is not None]))
    regiao_pos = total >= 0
    pares = sorted(((t, 0.0 if (c is None or (isinstance(c, float) and np.isnan(c))) else float(c))
                    for t, c in zip(tickers, contribs)),
                   key=lambda tc: tc[1], reverse=regiao_pos)
    cats, base, alta, baixa, tot = [], [], [], [], []
    cum = 0.0
    for tk, c in pares:
        cb, cum, ca = cum, cum + c, cum + c
        cats.append(tk)
        if regiao_pos:
            b, h = (cb if c >= 0 else ca), abs(c)
        else:
            b, h = max(cb, ca), -abs(c)
        base.append(b)
        alta.append(h if c >= 0 else None)
        baixa.append(h if c < 0 else None)
        tot.append(None)
    contrib_ord = [c for _, c in pares]      # contribuição real (com sinal), na ordem plotada
    cats.append(rotulo_total)
    base.append(0.0); alta.append(None); baixa.append(None); tot.append(total)
    return cats, base, alta, baixa, tot, contrib_ord


def _preenche_waterfall(chart, tickers, contribs, total):
    """Reescreve o gráfico de colunas do template como um waterfall
    (empilhado), escrevendo os caches direto nas séries."""
    cats, base, alta, baixa, tot, contrib_ord = _waterfall_arrays(tickers, contribs, total)
    n = len(cats)
    barChart = chart._chartSpace.find('.//' + qn('c:barChart'))
    barChart.find(qn('c:grouping')).set('val', 'stacked')
    ov = barChart.find(qn('c:overlap'))
    if ov is not None:
        ov.set('val', '100')
    for ser in barChart.findall(qn('c:ser')):
        barChart.remove(ser)

    cat_pts = ''.join(f'<c:pt idx="{i}"><c:v>{c}</c:v></c:pt>' for i, c in enumerate(cats))
    cat_xml = (f'<c:cat><c:strRef><c:f>Sheet1!$A$2:$A${n + 1}</c:f>'
               f'<c:strCache><c:ptCount val="{n}"/>{cat_pts}</c:strCache></c:strRef></c:cat>')
    nofill = '<c:spPr><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>'

    def _fill(rgb):
        return f'<c:spPr><a:solidFill><a:srgbClr val="{rgb}"/></a:solidFill></c:spPr>'

    def _fmt_sign(c):
        return f"{c * 100:+.1f}%".replace(".", ",")

    def _dlbls(labels):
        """labels: dict {idx: texto}. Rótulos de texto fixo (a altura da barra
        não é a contribuição, então não dá para usar showVal)."""
        if not labels:
            return ''
        itens = ''.join(
            f'<c:dLbl><c:idx val="{i}"/>'
            f'<c:tx><c:rich><a:bodyPr/><a:lstStyle/>'
            f'<a:p><a:pPr><a:defRPr sz="800"/></a:pPr>'
            f'<a:r><a:rPr lang="pt-BR" sz="800"/><a:t>{t}</a:t></a:r></a:p></c:rich></c:tx>'
            f'<c:dLblPos val="ctr"/>'
            f'<c:showLegendKey val="0"/><c:showVal val="0"/><c:showCatName val="0"/>'
            f'<c:showSerName val="0"/><c:showPercent val="0"/><c:showBubbleSize val="0"/></c:dLbl>'
            for i, t in sorted(labels.items()))
        return (f'{itens}<c:showLegendKey val="0"/><c:showVal val="0"/><c:showCatName val="0"/>'
                f'<c:showSerName val="0"/><c:showPercent val="0"/><c:showBubbleSize val="0"/>')

    def _ser(idx, nome, vals, sppr, labels=None):
        col = chr(ord('B') + idx)
        pts = ''.join('' if (v is None or (isinstance(v, float) and np.isnan(v)))
                      else f'<c:pt idx="{i}"><c:v>{float(v)}</c:v></c:pt>'
                      for i, v in enumerate(vals))
        dl = _dlbls(labels or {})
        dlbls_xml = f'<c:dLbls>{dl}</c:dLbls>' if dl else ''
        return (f'<c:ser xmlns:c="{_C_NS}" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<c:idx val="{idx}"/><c:order val="{idx}"/>'
                f'<c:tx><c:strRef><c:f>Sheet1!${col}$1</c:f><c:strCache><c:ptCount val="1"/>'
                f'<c:pt idx="0"><c:v>{nome}</c:v></c:pt></c:strCache></c:strRef></c:tx>'
                f'{sppr}{dlbls_xml}{cat_xml}'
                f'<c:val><c:numRef><c:f>Sheet1!${col}$2:${col}${n + 1}</c:f>'
                f'<c:numCache><c:formatCode>0.0%</c:formatCode>'
                f'<c:ptCount val="{n}"/>{pts}</c:numCache></c:numRef></c:val></c:ser>')

    # rótulos: contribuição real (com sinal) nas barras de alta/baixa; total na última
    lbl_alta = {i: _fmt_sign(contrib_ord[i]) for i in range(n - 1) if alta[i] is not None}
    lbl_baixa = {i: _fmt_sign(contrib_ord[i]) for i in range(n - 1) if baixa[i] is not None}
    lbl_total = {n - 1: _fmt_sign(tot[-1])}

    series = [
        _ser(0, 'Base', base, nofill),
        _ser(1, 'Contribuição (alta)', alta, _fill(COR_ALTA), lbl_alta),
        _ser(2, 'Contribuição (baixa)', baixa, _fill(COR_BAIXA), lbl_baixa),
        _ser(3, 'Carteira (total)', tot, _fill(COR_TOTAL), lbl_total),
    ]
    anchor = barChart.find(qn('c:dLbls'))
    if anchor is None:
        anchor = barChart.find(qn('c:gapWidth'))
    for s in series:
        anchor.addprevious(_parse_xml_c(s))


def _pc_grafico_decomposicao(shape, df_dec):
    """Waterfall: contribuição de cada papel + barra final com o retorno da
    carteira. A linha 'Carteira (total)' de df_dec vira o total."""
    total_row = df_dec[df_dec['Ticker'].astype(str).str.strip() == '']
    total = float(total_row['Contribuição'].iloc[0]) if not total_row.empty else None
    d = df_dec[df_dec['Ticker'].astype(str).str.strip() != ''].copy()
    _preenche_waterfall(shape.chart, d['Ticker'].tolist(),
                        d['Contribuição'].tolist(), total)


def _pc_datas(prs, ano_ref, mes_ref):
    """Atualiza os rótulos de mês. mês_carteira = mês_ref + 1."""
    ano_c, mes_c = _prox_mes(ano_ref, mes_ref)
    ref_nome = MESES_EXT_PT[mes_ref - 1].capitalize()
    cart_nome = MESES_EXT_PT[mes_c - 1].capitalize()
    alvos = {
        (1, 'object 7'): _mes_ext(ano_c, mes_c),
        (2, 'object 7'): _mes_ext(ano_ref, mes_ref),
        (2, 'CaixaDeTexto 8'): f"Alterações na Carteira – {cart_nome}",
        (2, 'CaixaDeTexto 24'): f"Decomposição do retorno da carteira ({ref_nome}/{str(ano_ref)[2:]})",
        (2, 'CaixaDeTexto 9'): f"Composição da Carteira – {cart_nome} de {str(ano_c)[2:]}",
        (3, 'object 7'): f"1 de {cart_nome} de {ano_c}",
    }
    for si, slide in enumerate(prs.slides, 1):
        for shp in slide.shapes:
            novo = alvos.get((si, shp.name))
            if novo is not None and shp.has_text_frame:
                _set_para_text(shp.text_frame.paragraphs[0], novo)


def atualizar_prestacao_contas(caminho_template, caminho_saida, df_port,
                               composition, portfolio, ano_ref=None, mes_ref=None):
    """Gera a Prestação de Contas. Por padrão, o mês de referência é o mês da
    última data de df_port (mês fechado); a composição e as 'alterações' são
    do mês seguinte. Passe ano_ref/mes_ref para fixar o mês de referência."""
    if ano_ref is None or mes_ref is None:
        ult = df_port.dropna(how='all').index.max()
        ano_ref, mes_ref = ult.year, ult.month

    df_comp = montar_df_composicao_compacta(composition, idioma='PT')
    df_dec = decomposicao_retorno(composition, ano_ref, mes_ref)

    prs = Presentation(caminho_template)
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_table:
                h = shp.table.rows[0].cells[0].text.strip()
                if h.startswith('Desempenho'):
                    _pc_tabela_resumo(shp.table, df_port, ano_ref, mes_ref)
                elif h.startswith('Setor'):
                    _pc_tabela_composicao(shp.table, df_comp)
            if shp.has_chart:
                if shp.chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED:
                    _pc_grafico_decomposicao(shp, df_dec)
                else:
                    _atualiza_grafico(shp, df_port)   # gráfico de linha base 100
    _pc_datas(prs, ano_ref, mes_ref)
    prs.save(caminho_saida)
    print(f"Prestação de Contas atualizada: {caminho_saida}")


# ------------------------------------------------------------------
# Caminhos: template e saída da Prestação de Contas por carteira.
# Só as entradas com 'template' existente são processadas.
# ------------------------------------------------------------------
prestacao_config = {
    'Carteira - TOP Ações XP': {
        'template': r"\\xpdocs\Research\Equities\Estrategia\Carteiras\Carteiras de Ações XP\Templates\Prestação de Contas - Top Ações.pptx",
        'saida':    r"\\xpdocs\Research\Equities\Estrategia\Carteiras\Carteiras de Ações XP\Prestação de Contas\Prestação de Contas - Top Ações.pptx",
    },
    # Preencha quando tiver os templates das demais carteiras:
    # 'Carteira - TOP DIVIDENDOS XP':  {'template': r"...", 'saida': r"..."},
    # 'Carteira - TOP SMALL CAPS XP':  {'template': r"...", 'saida': r"..."},
}

for portfolio, cfg in prestacao_config.items():
    template = cfg.get('template', '')
    if not template or not os.path.exists(template):
        print(f"[PULADO] template de prestação não encontrado para {portfolio}: {template}")
        continue
    atualizar_prestacao_contas(
        caminho_template=template,
        caminho_saida=cfg['saida'],
        df_port=_df_para_lamina(portfolio),
        composition=composition_dict[portfolio],
        portfolio=portfolio,
    )
