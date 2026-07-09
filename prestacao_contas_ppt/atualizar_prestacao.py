#!/usr/bin/env python3
"""Preenche o PPT de "Prestação de Contas" das carteiras a partir dos dados
já calculados pelo pipeline (portfolio_automation.py), preservando toda a
formatação do template.

Este módulo é **auto-contido** e recebe os dados como DataFrames — não acessa
rede nem depende do pipeline. O pipeline calcula os três insumos e chama
`atualizar_prestacao_contas`:

    df_port  — base 100 diária, colunas [carteira, benchmarks], index = datas
               (mesmo objeto usado na lâmina: _df_para_lamina(portfolio))
    df_comp  — composição compacta: colunas Setor, Companhia, Ticker, Peso
               (saída de montar_df_composicao_compacta(composition, 'PT'))
    df_dec   — decomposição de retorno do mês de referência: colunas
               Companhia, Ticker, Setor, Peso, 'Retorno no mês', 'Contribuição'
               (saída de decomposicao_retorno(composition, ano_ref, mes_ref);
                a linha 'Carteira (total)' é ignorada no gráfico)

O que é atualizado no template:
  * Tabela 3x4 "Desempenho" (mês de referência / acumulado no ano / 12 meses);
  * Gráfico de linha base 100 (carteira vs. benchmark);
  * Gráfico de colunas com a decomposição do retorno por papel;
  * Tabela de composição (Setor | Companhia | Ticker | Peso), com o Setor
    mesclado por grupo, crescendo/encolhendo conforme o nº de papéis;
  * As datas dos slides (mês de referência e mês da carteira = referência + 1).

Continuam manuais: os comentários do analista e o texto de "Alterações na
Carteira" (conteúdo editorial).

Requisitos: pip install python-pptx pandas numpy
"""
import datetime as dt
from copy import deepcopy

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn

C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
MESES_EXT_PT = ['janeiro', 'fevereiro', 'março', 'abril', 'maio', 'junho',
                'julho', 'agosto', 'setembro', 'outubro', 'novembro', 'dezembro']

# Mapeia o rótulo da linha/série do template para o nome da coluna em df_port
BENCH_LBL = {'Ibovespa': 'Ibovespa', 'SMLL': 'SMLL', 'ISEE': 'ISEE'}


# ----------------------------- formatação --------------------------------
def _fmt_pct_br(x, dec=1, dash='-'):
    if x is None or (isinstance(x, float) and np.isnan(x)):
        return dash
    return f"{x * 100:.{dec}f}%".replace(".", ",")


# --------------------------- retornos (base 100) -------------------------
def _ret_mes(serie, ano, mes):
    s = serie.dropna()
    ini = pd.Timestamp(ano, mes, 1) - pd.Timedelta(days=1)
    fim = pd.Timestamp(ano, mes, 1) + pd.offsets.MonthEnd(0)
    a, b = s[s.index <= ini], s[s.index <= fim]
    return np.nan if a.empty or b.empty else b.iloc[-1] / a.iloc[-1] - 1


def _ret_ano(serie, ano):
    s = serie.dropna()
    ini = pd.Timestamp(ano, 1, 1) - pd.Timedelta(days=1)
    fim = pd.Timestamp(ano, 12, 31)
    a, b = s[s.index <= ini], s[s.index <= fim]
    v0 = a.iloc[-1] if not a.empty else s.iloc[0]
    return np.nan if b.empty else b.iloc[-1] / v0 - 1


def _ret_periodo(serie, data_ini, data_fim):
    s = serie.dropna()
    a, b = s[s.index <= data_ini], s[s.index <= data_fim]
    v0 = a.iloc[-1] if not a.empty else s.iloc[0]
    return np.nan if b.empty else b.iloc[-1] / v0 - 1


# ------------------------------- texto -----------------------------------
def _set_para_text(p, value):
    """Troca o texto de um parágrafo preservando a formatação do 1º run."""
    runs = p.runs
    if runs:
        runs[0].text = str(value)
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p.add_run().text = str(value)


def _set_cell_text(cell, value):
    _set_para_text(cell.text_frame.paragraphs[0], value)


# ------------------------------ tabelas ----------------------------------
def _ajusta_linhas(table, n):
    """Deixa a tabela com n linhas de dados (fora o cabeçalho)."""
    tbl = table._tbl
    while len(tbl.tr_lst) - 1 < n:
        tbl.append(deepcopy(tbl.tr_lst[-1]))
    while len(tbl.tr_lst) - 1 > n:
        tbl.remove(tbl.tr_lst[-1])


def _unmerge_all(table):
    for r in range(len(table.rows)):
        for c in range(len(table.columns)):
            if table.cell(r, c).is_merge_origin:
                table.cell(r, c).split()


def _mes_ext(ano, mes):
    return f"{MESES_EXT_PT[mes - 1].capitalize()} de {ano}"


def _prox_mes(ano, mes):
    return (ano + (mes // 12), (mes % 12) + 1)


def _pc_tabela_resumo(table, df_port, ano_ref, mes_ref):
    """Tabela 3x4: Desempenho | <mês> | Acumulado <ano> | Últimos 12 meses."""
    nome_cart = df_port.columns[0]
    data_fim = df_port.dropna(how='all').index.max()
    ini_12m = data_fim - pd.DateOffset(months=12)
    hdr = table.rows[0].cells
    _set_cell_text(hdr[1], f"{MESES_EXT_PT[mes_ref - 1].capitalize()} {ano_ref}")
    _set_cell_text(hdr[2], f"Acumulado {ano_ref}")
    for row in list(table.rows)[1:]:
        rot = row.cells[0].text.strip()
        if 'carteira' in rot.lower():
            serie = nome_cart
        else:
            serie = next((c for c in df_port.columns
                          if c == rot or BENCH_LBL.get(c) == rot), None)
        if serie is None or serie not in df_port.columns:
            continue
        s = df_port[serie]
        _set_cell_text(row.cells[1], _fmt_pct_br(_ret_mes(s, ano_ref, mes_ref)))
        _set_cell_text(row.cells[2], _fmt_pct_br(_ret_ano(s, ano_ref)))
        _set_cell_text(row.cells[3], _fmt_pct_br(_ret_periodo(s, ini_12m, data_fim)))


def _pc_tabela_composicao(table, df_comp):
    """Setor | Companhia | Ticker | Peso, com o Setor mesclado por grupo."""
    _unmerge_all(table)
    n = len(df_comp)
    _ajusta_linhas(table, n)
    for i, r in enumerate(df_comp.itertuples(index=False), start=1):
        _set_cell_text(table.cell(i, 0), str(r.Setor))
        _set_cell_text(table.cell(i, 1), str(r.Companhia))
        _set_cell_text(table.cell(i, 2), str(r.Ticker))
        _set_cell_text(table.cell(i, 3), _fmt_pct_br(r.Peso))
    # mescla setores iguais consecutivos (esvazia continuações e mescla)
    r = 1
    while r <= n:
        start, val = r, table.cell(r, 0).text.strip()
        r += 1
        while r <= n and table.cell(r, 0).text.strip() == val:
            r += 1
        if r - 1 > start:
            for k in range(start + 1, r):
                _set_cell_text(table.cell(k, 0), "")
            table.cell(start, 0).merge(table.cell(r - 1, 0))


# ------------------------------ gráficos ---------------------------------
def _nome_serie_xml(serie_element):
    tx = serie_element.find(qn('c:tx'))
    if tx is None:
        return None
    vals = [v.text for v in tx.iter(qn('c:v')) if v.text]
    return vals[0] if vals else None


def _atualiza_grafico_linha(shape, df_port):
    """Substitui os dados do gráfico de linha (base 100), preservando estilo
    e os nomes originais das séries."""
    chart = shape.chart
    nome_cart = df_port.columns[0]
    cols = [nome_cart] + [c for c in df_port.columns if c != nome_cart]
    nomes = [_nome_serie_xml(s._element) for s in chart.series]
    cd = CategoryChartData(number_format=r'[$-416]mmm\-yy;@')
    cd.categories = [pd.Timestamp(d).to_pydatetime() for d in df_port.index]
    for i, col in enumerate(cols):
        nome = nomes[i] if i < len(nomes) and nomes[i] else col
        cd.add_series(nome, [None if pd.isna(v) else float(v) for v in df_port[col]])
    chart.replace_data(cd)


# Cores das barras do waterfall (alta / baixa / total). Amarelo XP no total.
COR_ALTA, COR_BAIXA, COR_TOTAL = '70AD47', 'C0504D', 'FFBC00'


def _waterfall_arrays(tickers, contribs, total, rotulo_total='Carteira'):
    """Monta os vetores de um waterfall com colunas empilhadas.

    Truque: uma série 'base' invisível flutua cada barra até o acumulado; as
    séries de alta/baixa desenham a variação e a última barra é o total.
    O problema clássico é o acumulado cruzar o zero (a base deixaria de ser
    um simples deslocamento). Solução: ordenar as contribuições para o
    caminho ficar sempre do mesmo lado do zero — positivas primeiro quando o
    total é >= 0 (caminho >= 0), negativas primeiro caso contrário — e
    calcular base/altura conforme a região.
    """
    total = float(total) if total is not None and not (isinstance(total, float) and np.isnan(total)) \
        else float(np.nansum([c for c in contribs if c is not None]))
    regiao_pos = total >= 0
    pares = sorted(((t, 0.0 if (c is None or (isinstance(c, float) and np.isnan(c))) else float(c))
                    for t, c in zip(tickers, contribs)),
                   key=lambda tc: tc[1], reverse=regiao_pos)
    cats, base, alta, baixa, tot = [], [], [], [], []
    cum = 0.0
    for tk, c in pares:
        cb, cum, ca = cum, cum + c, cum + c
        cats.append(tk)
        if regiao_pos:                     # caminho >= 0: base e alturas >= 0
            b, h = (cb if c >= 0 else ca), abs(c)
        else:                              # caminho <= 0: base e alturas <= 0
            b, h = max(cb, ca), -abs(c)
        base.append(b)
        alta.append(h if c >= 0 else None)
        baixa.append(h if c < 0 else None)
        tot.append(None)
    contrib_ord = [c for _, c in pares]      # contribuição real (com sinal), na ordem plotada
    cats.append(rotulo_total)
    base.append(0.0); alta.append(None); baixa.append(None); tot.append(total)
    return cats, base, alta, baixa, tot, contrib_ord


def _preenche_waterfall(chart, tickers, contribs, total):
    """Reescreve o gráfico de colunas do template como um waterfall
    (empilhado), escrevendo os caches direto nas séries."""
    cats, base, alta, baixa, tot, contrib_ord = _waterfall_arrays(tickers, contribs, total)
    n = len(cats)
    barChart = chart._chartSpace.find('.//' + qn('c:barChart'))
    barChart.find(qn('c:grouping')).set('val', 'stacked')
    ov = barChart.find(qn('c:overlap'))
    if ov is not None:
        ov.set('val', '100')
    for ser in barChart.findall(qn('c:ser')):
        barChart.remove(ser)

    cat_pts = ''.join(f'<c:pt idx="{i}"><c:v>{c}</c:v></c:pt>' for i, c in enumerate(cats))
    cat_xml = (f'<c:cat><c:strRef><c:f>Sheet1!$A$2:$A${n + 1}</c:f>'
               f'<c:strCache><c:ptCount val="{n}"/>{cat_pts}</c:strCache></c:strRef></c:cat>')
    nofill = '<c:spPr><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>'

    def _fill(rgb):
        return f'<c:spPr><a:solidFill><a:srgbClr val="{rgb}"/></a:solidFill></c:spPr>'

    def _fmt_sign(c):
        return f"{c * 100:+.1f}%".replace(".", ",")

    def _dlbls(labels):
        """labels: dict {idx: texto}. Rótulos de texto fixo (a altura da barra
        não é a contribuição, então não dá para usar showVal)."""
        if not labels:
            return ''
        itens = ''.join(
            f'<c:dLbl><c:idx val="{i}"/>'
            f'<c:tx><c:rich><a:bodyPr/><a:lstStyle/>'
            f'<a:p><a:pPr><a:defRPr sz="800"/></a:pPr>'
            f'<a:r><a:rPr lang="pt-BR" sz="800"/><a:t>{t}</a:t></a:r></a:p></c:rich></c:tx>'
            f'<c:showLegendKey val="0"/><c:showVal val="0"/><c:showCatName val="0"/>'
            f'<c:showSerName val="0"/><c:showPercent val="0"/><c:showBubbleSize val="0"/></c:dLbl>'
            for i, t in sorted(labels.items()))
        return (f'{itens}<c:showLegendKey val="0"/><c:showVal val="0"/><c:showCatName val="0"/>'
                f'<c:showSerName val="0"/><c:showPercent val="0"/><c:showBubbleSize val="0"/>')

    def _ser(idx, nome, vals, sppr, labels=None):
        col = chr(ord('B') + idx)
        pts = ''.join('' if (v is None or (isinstance(v, float) and np.isnan(v)))
                      else f'<c:pt idx="{i}"><c:v>{float(v)}</c:v></c:pt>'
                      for i, v in enumerate(vals))
        dl = _dlbls(labels or {})
        dlbls_xml = f'<c:dLbls>{dl}</c:dLbls>' if dl else ''
        return (f'<c:ser xmlns:c="{C_NS}" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<c:idx val="{idx}"/><c:order val="{idx}"/>'
                f'<c:tx><c:strRef><c:f>Sheet1!${col}$1</c:f><c:strCache><c:ptCount val="1"/>'
                f'<c:pt idx="0"><c:v>{nome}</c:v></c:pt></c:strCache></c:strRef></c:tx>'
                f'{sppr}{dlbls_xml}{cat_xml}'
                f'<c:val><c:numRef><c:f>Sheet1!${col}$2:${col}${n + 1}</c:f>'
                f'<c:numCache><c:formatCode>0.0%</c:formatCode>'
                f'<c:ptCount val="{n}"/>{pts}</c:numCache></c:numRef></c:val></c:ser>')

    # rótulos: contribuição real (com sinal) nas barras de alta/baixa; total na última
    lbl_alta = {i: _fmt_sign(contrib_ord[i]) for i in range(n - 1) if alta[i] is not None}
    lbl_baixa = {i: _fmt_sign(contrib_ord[i]) for i in range(n - 1) if baixa[i] is not None}
    lbl_total = {n - 1: _fmt_sign(tot[-1])}

    series = [
        _ser(0, 'Base', base, nofill),
        _ser(1, 'Contribuição (alta)', alta, _fill(COR_ALTA), lbl_alta),
        _ser(2, 'Contribuição (baixa)', baixa, _fill(COR_BAIXA), lbl_baixa),
        _ser(3, 'Carteira (total)', tot, _fill(COR_TOTAL), lbl_total),
    ]
    anchor = barChart.find(qn('c:dLbls'))
    if anchor is None:
        anchor = barChart.find(qn('c:gapWidth'))
    for s in series:
        anchor.addprevious(parse_xml(s))


def _pc_grafico_decomposicao(shape, df_dec):
    """Waterfall: contribuição de cada papel + barra final com o retorno da
    carteira. A linha 'Carteira (total)' de df_dec vira o total."""
    total_row = df_dec[df_dec['Ticker'].astype(str).str.strip() == '']
    total = float(total_row['Contribuição'].iloc[0]) if not total_row.empty else None
    d = df_dec[df_dec['Ticker'].astype(str).str.strip() != ''].copy()
    _preenche_waterfall(shape.chart, d['Ticker'].tolist(),
                        d['Contribuição'].tolist(), total)


# -------------------------------- datas ----------------------------------
def _pc_datas(prs, ano_ref, mes_ref):
    """Atualiza os rótulos de mês. mês_carteira = mês_ref + 1."""
    ano_c, mes_c = _prox_mes(ano_ref, mes_ref)
    ref_nome = MESES_EXT_PT[mes_ref - 1].capitalize()
    cart_nome = MESES_EXT_PT[mes_c - 1].capitalize()
    # (slide, nome do shape) -> novo texto do 1º parágrafo
    alvos = {
        (1, 'object 7'): _mes_ext(ano_c, mes_c),
        (2, 'object 7'): _mes_ext(ano_ref, mes_ref),
        (2, 'CaixaDeTexto 8'): f"Alterações na Carteira – {cart_nome}",
        (2, 'CaixaDeTexto 24'): f"Decomposição do retorno da carteira ({ref_nome}/{str(ano_ref)[2:]})",
        (2, 'CaixaDeTexto 9'): f"Composição da Carteira – {cart_nome} de {str(ano_c)[2:]}",
        (3, 'object 7'): f"1 de {cart_nome} de {ano_c}",
    }
    for si, slide in enumerate(prs.slides, 1):
        for shp in slide.shapes:
            novo = alvos.get((si, shp.name))
            if novo is not None and shp.has_text_frame:
                _set_para_text(shp.text_frame.paragraphs[0], novo)


# ----------------------------- orquestração ------------------------------
def atualizar_prestacao_contas(caminho_template, caminho_saida,
                               df_port, df_comp, df_dec, ano_ref, mes_ref):
    """Abre o template, preenche as tabelas/gráficos/datas e salva."""
    prs = Presentation(caminho_template)
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_table:
                h = shp.table.rows[0].cells[0].text.strip()
                if h.startswith('Desempenho'):
                    _pc_tabela_resumo(shp.table, df_port, ano_ref, mes_ref)
                elif h.startswith('Setor'):
                    _pc_tabela_composicao(shp.table, df_comp)
            if shp.has_chart:
                if shp.chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED:
                    _pc_grafico_decomposicao(shp, df_dec)
                else:
                    _atualiza_grafico_linha(shp, df_port)
    _pc_datas(prs, ano_ref, mes_ref)
    prs.save(caminho_saida)
    print(f"Prestação de Contas atualizada: {caminho_saida}")
